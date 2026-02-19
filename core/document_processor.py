"""Document processor — text extraction, MIME routing, chunking, OCR.

Handles all supported file types: PDF, DOCX, XLSX, PPTX, EPUB, images, plain text.
CPU-bound work is wrapped in asyncio.to_thread for non-blocking execution.
"""

from __future__ import annotations

import asyncio
import hashlib
import logging
import mimetypes
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

from core.config import DocumentConfig

logger = logging.getLogger(__name__)


@dataclass
class ExtractedContent:
    """Result of extracting text from a document."""

    text: str
    page_count: int = 1
    metadata: dict[str, Any] = field(default_factory=dict)
    images: list[bytes] = field(default_factory=list)
    mime_type: str = ""
    pages: list[str] = field(default_factory=list)  # per-page text for PDFs


@dataclass
class DocumentChunk:
    """A single chunk ready for embedding."""

    content: str
    chunk_index: int
    token_count: int
    page_number: int | None = None
    section_title: str = ""
    metadata: dict[str, Any] = field(default_factory=dict)


class DocumentProcessor:
    """Extracts text, chunks, and prepares documents for embedding or direct analysis."""

    def __init__(self, config: DocumentConfig) -> None:
        self._config = config
        self._ocr_engine: Any = None  # Lazy-loaded RapidOCR

    @staticmethod
    def estimate_tokens(text: str) -> int:
        """Token count estimate (~4 chars per token, matching core/indexer.py)."""
        return len(text) // 4

    @staticmethod
    def content_hash(file_path: Path) -> str:
        """SHA-256 hash of file contents."""
        h = hashlib.sha256()
        with open(file_path, "rb") as f:
            for chunk in iter(lambda: f.read(65536), b""):
                h.update(chunk)
        return h.hexdigest()

    def detect_mime_type(self, file_path: Path) -> str:
        """Detect MIME type using extension, with magic-byte fallback for PDFs."""
        mime, _ = mimetypes.guess_type(str(file_path))
        if mime:
            return mime

        # Check magic bytes for common types
        try:
            with open(file_path, "rb") as f:
                header = f.read(8)
            if header.startswith(b"%PDF"):
                return "application/pdf"
            if header[:4] == b"PK\x03\x04":
                # Could be docx, xlsx, pptx, epub — check extension
                suffix = file_path.suffix.lower()
                ext_map = {
                    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                    ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    ".epub": "application/epub+zip",
                }
                return ext_map.get(suffix, "application/zip")
            if header[:3] in (b"\xff\xd8\xff", ):
                return "image/jpeg"
            if header[:8] == b"\x89PNG\r\n\x1a\n":
                return "image/png"
        except OSError:
            pass

        return "application/octet-stream"

    def is_image(self, mime_type: str) -> bool:
        return mime_type.startswith("image/")

    async def extract_text(self, file_path: Path, mime_type: str | None = None) -> ExtractedContent:
        """Route to the correct extractor based on MIME type."""
        if mime_type is None:
            mime_type = self.detect_mime_type(file_path)

        extractors = {
            "application/pdf": self._extract_pdf,
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document": self._extract_docx,
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": self._extract_xlsx,
            "text/csv": self._extract_csv,
            "application/vnd.openxmlformats-officedocument.presentationml.presentation": self._extract_pptx,
            "application/epub+zip": self._extract_epub,
        }

        # Check exact match
        extractor = extractors.get(mime_type)
        if extractor:
            return await extractor(file_path)

        # Image types → OCR
        if self.is_image(mime_type):
            return await self._extract_image_ocr(file_path)

        # Text types → plain text
        if mime_type.startswith("text/"):
            return await self._extract_plain(file_path)

        # Fallback: try plain text
        logger.warning("Unknown MIME type %s for %s, trying plain text", mime_type, file_path.name)
        return await self._extract_plain(file_path)

    async def _extract_pdf(self, file_path: Path) -> ExtractedContent:
        """Extract text from PDF using pymupdf. Falls back to OCR for scanned pages."""
        def _extract() -> ExtractedContent:
            import pymupdf

            doc = pymupdf.open(str(file_path))
            pages: list[str] = []
            all_text: list[str] = []
            metadata = {}

            if doc.metadata:
                for key in ("title", "author", "subject"):
                    val = doc.metadata.get(key)
                    if val:
                        metadata[key] = val

            for page in doc:
                text = page.get_text()
                pages.append(text)
                all_text.append(text)

            doc.close()
            full_text = "\n\n".join(all_text)

            return ExtractedContent(
                text=full_text,
                page_count=len(pages),
                metadata=metadata,
                mime_type="application/pdf",
                pages=pages,
            )

        result = await asyncio.to_thread(_extract)

        # If very little text extracted, try OCR on the whole file
        if self._config.ocr_enabled and self.estimate_tokens(result.text) < 50 and result.page_count > 0:
            logger.info("PDF appears scanned, attempting OCR on %s", file_path.name)
            ocr_result = await self._ocr_pdf(file_path)
            if self.estimate_tokens(ocr_result.text) > self.estimate_tokens(result.text):
                return ocr_result

        return result

    async def _ocr_pdf(self, file_path: Path) -> ExtractedContent:
        """OCR a scanned PDF by rendering pages to images."""
        def _ocr() -> ExtractedContent:
            import pymupdf

            engine = self._get_ocr_engine()
            doc = pymupdf.open(str(file_path))
            pages: list[str] = []

            for page in doc:
                pix = page.get_pixmap(dpi=200)
                img_bytes = pix.tobytes("png")
                result, _ = engine(img_bytes)
                page_text = ""
                if result:
                    page_text = "\n".join(line[1] for line in result)
                pages.append(page_text)

            doc.close()
            return ExtractedContent(
                text="\n\n".join(pages),
                page_count=len(pages),
                mime_type="application/pdf",
                pages=pages,
            )

        return await asyncio.to_thread(_ocr)

    async def _extract_docx(self, file_path: Path) -> ExtractedContent:
        def _extract() -> ExtractedContent:
            from docx import Document

            doc = Document(str(file_path))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            metadata = {}
            if doc.core_properties.title:
                metadata["title"] = doc.core_properties.title
            if doc.core_properties.author:
                metadata["author"] = doc.core_properties.author

            return ExtractedContent(
                text="\n\n".join(paragraphs),
                page_count=max(1, len(paragraphs) // 30),  # approximate
                metadata=metadata,
                mime_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            )

        return await asyncio.to_thread(_extract)

    async def _extract_xlsx(self, file_path: Path) -> ExtractedContent:
        def _extract() -> ExtractedContent:
            from openpyxl import load_workbook

            wb = load_workbook(str(file_path), read_only=True, data_only=True)
            sheets: list[str] = []

            for ws in wb.worksheets:
                rows: list[str] = []
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) if c is not None else "" for c in row]
                    if any(cells):
                        rows.append("\t".join(cells))
                if rows:
                    sheets.append(f"## Sheet: {ws.title}\n" + "\n".join(rows))

            wb.close()
            return ExtractedContent(
                text="\n\n".join(sheets),
                page_count=len(sheets),
                metadata={"sheet_count": len(sheets)},
                mime_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            )

        return await asyncio.to_thread(_extract)

    async def _extract_csv(self, file_path: Path) -> ExtractedContent:
        def _extract() -> ExtractedContent:
            import csv

            with open(file_path, newline="", encoding="utf-8", errors="replace") as f:
                reader = csv.reader(f)
                rows = ["\t".join(row) for row in reader]

            return ExtractedContent(
                text="\n".join(rows),
                page_count=1,
                mime_type="text/csv",
            )

        return await asyncio.to_thread(_extract)

    async def _extract_pptx(self, file_path: Path) -> ExtractedContent:
        def _extract() -> ExtractedContent:
            from pptx import Presentation

            prs = Presentation(str(file_path))
            slides: list[str] = []

            for i, slide in enumerate(prs.slides, 1):
                parts: list[str] = [f"## Slide {i}"]
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text = shape.text_frame.text.strip()
                        if text:
                            parts.append(text)
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        parts.append(f"[Speaker notes: {notes}]")
                slides.append("\n".join(parts))

            return ExtractedContent(
                text="\n\n".join(slides),
                page_count=len(slides),
                metadata={"slide_count": len(slides)},
                mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )

        return await asyncio.to_thread(_extract)

    async def _extract_epub(self, file_path: Path) -> ExtractedContent:
        def _extract() -> ExtractedContent:
            from html.parser import HTMLParser

            import ebooklib
            from ebooklib import epub

            book = epub.read_epub(str(file_path), options={"ignore_ncx": True})
            chapters: list[str] = []

            class _TextExtractor(HTMLParser):
                def __init__(self) -> None:
                    super().__init__()
                    self.parts: list[str] = []
                def handle_data(self, data: str) -> None:
                    text = data.strip()
                    if text:
                        self.parts.append(text)

            for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
                parser = _TextExtractor()
                parser.feed(item.get_content().decode("utf-8", errors="replace"))
                if parser.parts:
                    chapters.append("\n".join(parser.parts))

            metadata = {}
            title = book.get_metadata("DC", "title")
            if title:
                metadata["title"] = title[0][0]

            return ExtractedContent(
                text="\n\n".join(chapters),
                page_count=len(chapters),
                metadata=metadata,
                mime_type="application/epub+zip",
            )

        return await asyncio.to_thread(_extract)

    async def _extract_plain(self, file_path: Path) -> ExtractedContent:
        def _extract() -> ExtractedContent:
            from charset_normalizer import from_path

            result = from_path(file_path)
            best = result.best()
            text = str(best) if best else file_path.read_text(errors="replace")

            return ExtractedContent(
                text=text,
                page_count=1,
                mime_type="text/plain",
            )

        return await asyncio.to_thread(_extract)

    async def _extract_image_ocr(self, file_path: Path) -> ExtractedContent:
        """OCR an image using rapidocr-onnxruntime."""
        if not self._config.ocr_enabled:
            return ExtractedContent(
                text="[Image — OCR disabled]",
                mime_type=self.detect_mime_type(file_path),
            )

        def _ocr() -> ExtractedContent:
            engine = self._get_ocr_engine()
            img_bytes = file_path.read_bytes()
            result, _ = engine(img_bytes)
            text = ""
            if result:
                text = "\n".join(line[1] for line in result)
            return ExtractedContent(
                text=text or "[No text detected in image]",
                mime_type=self.detect_mime_type(file_path),
            )

        return await asyncio.to_thread(_ocr)

    def _get_ocr_engine(self) -> Any:
        """Lazy-load the OCR engine."""
        if self._ocr_engine is None:
            from rapidocr_onnxruntime import RapidOCR
            self._ocr_engine = RapidOCR()
        return self._ocr_engine

    def should_use_rag(self, content: ExtractedContent) -> bool:
        """Return True if document exceeds context_threshold_tokens."""
        return self.estimate_tokens(content.text) > self._config.context_threshold_tokens

    def chunk_document(
        self, content: ExtractedContent, source_filename: str
    ) -> list[DocumentChunk]:
        """Split extracted content into chunks based on document type."""
        # PDF with pages: chunk by pages
        if content.pages and len(content.pages) > 1:
            return self._chunk_by_pages(content.pages)

        # Try structural chunking by headings
        if "\n## " in content.text or "\n# " in content.text:
            return self._chunk_by_sections(content.text)

        # Fallback: sliding window
        return self._chunk_by_size(content.text)

    def _chunk_by_pages(self, pages: list[str]) -> list[DocumentChunk]:
        """Split by page boundaries, merging small pages."""
        chunks: list[DocumentChunk] = []
        buffer = ""
        buffer_start_page = 1
        max_tokens = self._config.chunk_size_tokens

        for i, page_text in enumerate(pages, 1):
            page_text = page_text.strip()
            if not page_text:
                continue

            combined = (buffer + "\n\n" + page_text).strip() if buffer else page_text
            if self.estimate_tokens(combined) <= max_tokens:
                buffer = combined
                if not chunks and not buffer.strip():
                    buffer_start_page = i
            else:
                # Flush buffer as chunk
                if buffer.strip():
                    chunks.append(DocumentChunk(
                        content=buffer.strip(),
                        chunk_index=len(chunks),
                        token_count=self.estimate_tokens(buffer),
                        page_number=buffer_start_page,
                    ))
                buffer = page_text
                buffer_start_page = i

        # Flush remaining
        if buffer.strip():
            chunks.append(DocumentChunk(
                content=buffer.strip(),
                chunk_index=len(chunks),
                token_count=self.estimate_tokens(buffer),
                page_number=buffer_start_page,
            ))

        return chunks

    def _chunk_by_sections(self, text: str) -> list[DocumentChunk]:
        """Split by heading markers."""
        import re

        sections: list[tuple[str, str]] = []  # (title, content)
        parts = re.split(r"(?=\n#{1,3}\s)", text)

        for part in parts:
            part = part.strip()
            if not part:
                continue
            # Extract heading
            lines = part.split("\n", 1)
            title = lines[0].lstrip("#").strip() if lines[0].startswith("#") else ""
            body = lines[1].strip() if len(lines) > 1 else part
            sections.append((title, body if body else part))

        chunks: list[DocumentChunk] = []
        max_tokens = self._config.chunk_size_tokens

        for title, body in sections:
            if self.estimate_tokens(body) <= max_tokens:
                chunks.append(DocumentChunk(
                    content=body,
                    chunk_index=len(chunks),
                    token_count=self.estimate_tokens(body),
                    section_title=title,
                ))
            else:
                # Section too large — split by size
                for sub in self._chunk_by_size(body):
                    sub.section_title = title
                    sub.chunk_index = len(chunks)
                    chunks.append(sub)

        return chunks

    def _chunk_by_size(self, text: str) -> list[DocumentChunk]:
        """Sliding window chunking with overlap."""
        max_tokens = self._config.chunk_size_tokens
        overlap_tokens = self._config.chunk_overlap_tokens
        max_chars = max_tokens * 4
        overlap_chars = overlap_tokens * 4

        chunks: list[DocumentChunk] = []
        start = 0

        while start < len(text):
            end = start + max_chars
            # Try to break at a paragraph boundary
            if end < len(text):
                newline_pos = text.rfind("\n\n", start + max_chars // 2, end + 200)
                if newline_pos > start:
                    end = newline_pos

            chunk_text = text[start:end].strip()
            if chunk_text:
                chunks.append(DocumentChunk(
                    content=chunk_text,
                    chunk_index=len(chunks),
                    token_count=self.estimate_tokens(chunk_text),
                ))

            start = end - overlap_chars
            if start <= (end - max_chars):
                start = end  # Prevent infinite loop

        return chunks
