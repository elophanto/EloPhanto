"""Executive deck rendering for the competitive-intelligence organ.

The scorecard workbook is for the analyst and the board report is for the
reader; this is for the room. ~16 slides in a fixed house style, built from
the *same* stored evidence as the other deliverables — nothing here is
computed differently, it is only said the way a room hears it.

Design doctrine (adapted from decks that ship to steering committees):

* **The room hears about the market, not the machinery.** The front half of
  the deck is competitors: an executive summary in findings / threats /
  watch-next columns, standings, us-versus-the-leader, one deep-dive slide
  per key competitor (observations → implications), their storefronts as
  photographed exhibits, and the moves this period. Evidence coverage and
  method exist — in the appendix, where an analyst looks for them.
* **One idea per slide, action titles.** Every heading is a sentence someone
  could disagree with ("High 5 leads a thin field"), never a label. The
  model writes titles and commentary from the factual record; the numbers
  themselves are computed, never generated.
* **Exhibits are captures, not mockups.** A storefront screenshot on a slide
  was taken by the browser through a state-verified network exit, and is
  filed in the evidence register beside the claims from that page.
* **Restraint.** White content slides, ink type, one short accent rule under
  each heading, dark bookends. En dashes, never em.
* **No internal bookkeeping.** Hashes, file paths, manifests, run IDs and
  checkpoint numbers never reach a slide — enforced by prompt *and* by a
  scrubber here, because a deck once shipped with a SHA-256 on it.

And the organ's honesty rules survive the trip onto slides, where they are
most easily lost:

* unscored is blank, never zero — no bar, an empty heatmap cell;
* provisional brands are listed beside the chart, never ranked in it;
* model-written narrative is labelled as such, and when no model is
  available the deck says *facts only* — "could not evaluate" is never
  dressed as "nothing to report".
"""

from __future__ import annotations

import logging
import re
from datetime import UTC, datetime
from pathlib import Path
from typing import Any

logger = logging.getLogger(__name__)

# ── house tokens ─────────────────────────────────────────────────────
_INK = "111827"  # near-black: headings, display numbers, dark canvases
_BODY = "4B5563"  # body copy
_MUTED = "9CA3AF"  # eyebrows, footers, labels
_HAIR = "E5E7EB"  # hairlines
_CARD = "F9FAFB"  # zebra rows / note cards
_GAP_BG = "F3F4F6"  # heatmap: not observed
_ACCENT = "D97706"  # amber: the accent rule, and *our* brand everywhere
_PEER = "64748B"  # slate: peer brands
_DARK_BODY = "D1D5DB"  # body copy on dark canvases
_SELF_ROW = "FDF6EC"  # our row in the heatmap
_WHITE = "FFFFFF"

_CLASS_LABEL = {
    "no_regret": "No-regret",
    "transition_requirement": "Transition requirement",
    "post_transition": "Post-transition",
    "monitor": "Monitor",
}
_CLASS_ORDER = ["no_regret", "transition_requirement", "post_transition", "monitor"]

# Internal-bookkeeping tokens that must never reach a slide. The narrative
# prompt bans them; this scrubs whatever slips through. Hex runs need at
# least one letter so ordinary numbers ("1000000") survive.
_HEX_RE = re.compile(r"\b(?=[0-9a-f]*[a-f])[0-9a-f]{7,64}\b", re.I)
_BOOKKEEPING_RE = re.compile(
    r"\b(sha-?\d{0,3}|checksum|manifest|freeze receipt|run.id|checkpoint \d+)\b",
    re.I,
)


def _clean(text: Any, cap: int = 300) -> str:
    """House copy rules: en dashes, no bookkeeping tokens, collapsed, capped."""
    s = str(text or "")
    s = s.replace("—", "–")
    s = _HEX_RE.sub("", s)
    s = _BOOKKEEPING_RE.sub("", s)
    s = re.sub(r"\s+", " ", s).strip(" -–,;")
    return s[:cap]


def _trim_words(text: str, cap: int) -> str:
    """Cap at a word boundary — a slide once shipped reading '…FLORI'."""
    if len(text) <= cap:
        return text
    cut = text[:cap]
    if " " in cut:
        cut = cut[: cut.rfind(" ")]
    return cut.rstrip(" ·-–,;")


def _fmt(v: float | None) -> str:
    return "—" if v is None else f"{v:.1f}"


# ── drawing primitives ───────────────────────────────────────────────


def _rgb(hexstr: str) -> Any:
    from pptx.dml.color import RGBColor

    return RGBColor.from_string(hexstr)


def _text(
    slide: Any,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    *,
    size: float = 14,
    bold: bool = False,
    italic: bool = False,
    color: str = _INK,
    align: str = "left",
    spacing: float | None = None,
    line: float | None = None,
    wrap: bool = True,
) -> Any:
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }[align]
    if line:
        p.line_spacing = line
    for r in p.runs:
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = _rgb(color)
        if spacing is not None:
            rpr = r._r.get_or_add_rPr()
            rpr.set("spc", str(int(spacing * 100)))
    return box


def _bullets(
    slide: Any,
    left: float,
    top: float,
    width: float,
    height: float,
    items: list[str],
    *,
    size: float = 15,
    color: str = _BODY,
    gap_pt: int = 10,
    cap: int = 220,
    accent_bullet: bool = True,
    max_items: int = 6,
) -> Any:
    from pptx.util import Inches, Pt

    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for item in items[:max_items]:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        if accent_bullet:
            r0 = p.add_run()
            r0.text = "•  "
            r0.font.size = Pt(size)
            r0.font.bold = True
            r0.font.color.rgb = _rgb(_ACCENT)
        r = p.add_run()
        r.text = _clean(item, cap)
        r.font.size = Pt(size)
        r.font.color.rgb = _rgb(color)
        p.space_after = Pt(gap_pt)
    return box


def _rule(slide: Any, y: float, *, x: float = 0.7, w: float = 0.7) -> None:
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt

    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Pt(2.6))
    ln.fill.solid()
    ln.fill.fore_color.rgb = _rgb(_ACCENT)
    ln.line.fill.background()
    ln.shadow.inherit = False


def _eyebrow(slide: Any, text: str, *, y: float, x: float = 0.7, color: str = _MUTED) -> None:
    # Width fits the remaining canvas — an eyebrow placed in a right-hand
    # column must not spill past the slide edge — and the text is cut at a
    # word boundary, never mid-word.
    width = max(1.0, 13.333 - x - 0.73)
    _text(
        slide,
        x,
        y,
        width,
        0.3,
        _trim_words(_clean(text, 90), 60).upper(),
        size=10.5,
        bold=True,
        color=color,
        spacing=3,
    )


def _footer(slide: Any, deck_title: str, page: int) -> None:
    _text(slide, 0.7, 7.08, 8.5, 0.3, _clean(deck_title, 70), size=8.5, color=_MUTED)
    _text(slide, 12.0, 7.08, 0.7, 0.3, str(page), size=8.5, color=_MUTED, align="right")


def _header(slide: Any, eyebrow: str, title: str, commentary: str = "") -> float:
    """Eyebrow, action title, accent rule, optional one-line takeaway.

    Returns the y where slide content should start.
    """
    _eyebrow(slide, eyebrow, y=0.42)
    _text(
        slide,
        0.7,
        0.72,
        11.9,
        0.85,
        _clean(title, 110),
        size=23,
        bold=True,
        line=1.08,
    )
    _rule(slide, 1.62)
    if commentary:
        _text(
            slide,
            0.7,
            1.78,
            11.9,
            0.4,
            _clean(commentary, 170),
            size=12.5,
            italic=True,
            color=_BODY,
        )
        return 2.3
    return 2.0


def _judgement_note(slide: Any, source: str, *, dark: bool = False) -> None:
    msg = (
        "Narrative and commentary written by the model from the factual record "
        "– verify before presenting."
        if source == "model"
        else "Facts only – no model was available, so no narrative judgement has been applied."
    )
    _text(slide, 0.7, 6.72, 11.9, 0.3, msg, size=9, color="6B7280" if dark else _MUTED)


def _notes(slide: Any, text: str) -> None:
    slide.notes_slide.notes_text_frame.text = text[:1000]


def _blank(prs: Any) -> Any:
    return prs.slides.add_slide(prs.slide_layouts[6])


def _dark(slide: Any) -> None:
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = _rgb(_INK)
    bg.line.fill.background()
    bg.shadow.inherit = False


def _picture(slide: Any, path: str, x: float, y: float, w: float, max_h: float) -> Any | None:
    """Place an image fitted to ``w`` wide, capped at ``max_h`` tall, with a
    hairline border. Returns the picture shape, or None if the file is
    unreadable — a missing exhibit never breaks the deck."""
    from pptx.util import Inches, Pt

    try:
        pic = slide.shapes.add_picture(path, Inches(x), Inches(y), width=Inches(w))
    except Exception as e:
        logger.debug("deck: could not place exhibit %s: %s", path, e)
        return None
    if pic.height > Inches(max_h):
        ratio = Inches(max_h) / pic.height
        pic.height = Inches(max_h)
        pic.width = int(pic.width * ratio)
    pic.line.color.rgb = _rgb(_HAIR)
    pic.line.width = Pt(1.0)
    pic.shadow.inherit = False
    return pic


# ── narrative fallback (no model) ────────────────────────────────────


def _sidebar(
    slide: Any,
    observations: list[str],
    implications: list[str],
    *,
    top: float,
    x: float = 9.1,
    w: float = 3.55,
    bottom: float = 6.55,
) -> float:
    """The right-hand reading panel every analytical slide carries: *Key
    observations* (what the chart shows) and *Key implications* (what it
    means for us). Executives read this column and skip the chart; the
    reference decks the customer benchmarks against carry it on every
    slide, and a slide without it makes the room guess. Returns the y where
    the panel ends. Draws nothing when both lists are empty."""
    obs = [str(o).strip() for o in observations if str(o).strip()][:4]
    imp = [str(i).strip() for i in implications if str(i).strip()][:3]
    if not obs and not imp:
        return top
    from pptx.util import Inches

    h = bottom - top
    card = slide.shapes.add_shape(
        1, Inches(x - 0.15), Inches(top - 0.1), Inches(w + 0.3), Inches(h)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = _rgb(_CARD)
    card.line.fill.background()
    card.shadow.inherit = False
    y = top + 0.05
    if obs:
        _text(slide, x, y, w, 0.3, "Key observations", size=11, bold=True, color=_INK)
        y += 0.36
        n_lines = sum(1 + len(o) // 48 for o in obs)
        block_h = min(0.24 * n_lines + 0.12 * len(obs), bottom - y - (1.6 if imp else 0.2))
        _bullets(
            slide,
            x,
            y,
            w,
            block_h,
            obs,
            size=9.5,
            color=_BODY,
            gap_pt=5,
            cap=150,
            accent_bullet=False,
            max_items=4,
        )
        y += block_h + 0.18
    if imp and y < bottom - 0.8:
        _text(slide, x, y, w, 0.3, "Key implications", size=11, bold=True, color=_ACCENT)
        y += 0.36
        _bullets(
            slide,
            x,
            y,
            w,
            max(0.5, bottom - y - 0.1),
            imp,
            size=9.5,
            color=_BODY,
            gap_pt=5,
            cap=150,
            accent_bullet=True,
            max_items=3,
        )
    return bottom


def _chip(slide: Any, x: float, y: float, n: int, *, color: str = _PEER) -> None:
    """A small numbered circle — the reference decks number each executive
    observation and repeat the number where the evidence lives."""
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    d = 0.26
    c = slide.shapes.add_shape(9, Inches(x), Inches(y), Inches(d), Inches(d))  # 9 = oval
    c.fill.solid()
    c.fill.fore_color.rgb = _rgb(color)
    c.line.fill.background()
    c.shadow.inherit = False
    tf = c.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = str(n)
    p.alignment = PP_ALIGN.CENTER
    for r in p.runs:
        r.font.size = Pt(8.5)
        r.font.bold = True
        r.font.color.rgb = _rgb(_WHITE)


def _slides_facts(
    card: dict[str, Any], offers: list[dict[str, Any]] | None = None
) -> dict[str, Any]:
    """Computed, model-free observations/implications per slide — the
    reading panel's fallback so a facts-only deck still tells the room what
    each slide shows. Numbers only, no claims about intent."""
    rows = card.get("rows", [])
    ranked = [r for r in rows if r.get("rank") is not None]
    us_rows = [r for r in rows if r.get("is_self")]
    us = us_rows[0] if us_rows else None
    leader = next((r for r in ranked if not r.get("is_self")), None)
    dims = card.get("dimensions", [])
    out: dict[str, dict[str, list[str]]] = {}

    st_obs: list[str] = []
    st_imp: list[str] = []
    if leader:
        st_obs.append(
            f"{leader['name']} leads the ranked field at {_fmt(leader['overall']['normalized_pct'])}."
        )
    if len(ranked) >= 3:
        spread = float(ranked[0]["overall"]["normalized_pct"]) - float(
            ranked[-1]["overall"]["normalized_pct"]
        )
        st_obs.append(f"{len(ranked)} brands ranked; {spread:.0f} points separate first from last.")
    for u in us_rows[:2]:
        if u.get("rank") is not None:
            st_obs.append(
                f"{u['name']} ranks #{u['rank']} at {_fmt(u['overall']['normalized_pct'])}."
            )
            if leader is not None:
                gap = float(leader["overall"]["normalized_pct"]) - float(
                    u["overall"]["normalized_pct"]
                )
                st_imp.append(f"{u['name']} sits {gap:.1f} points behind the leader.")
        elif u["overall"]["normalized_pct"] is not None:
            st_obs.append(f"{u['name']} is scored but not yet ranked.")
    if card.get("comparability_note"):
        st_imp.append("Ranks are withheld until the field is measured to comparable depth.")
    out["standings"] = {"observations": st_obs, "implications": st_imp}

    # dimensions: where we lead / trail
    lead, trail = [], []
    if us is not None:
        for d in dims:
            dn = d["name"]
            mine = us.get("dimensions", {}).get(dn, {}).get("score")
            scored = [
                float(r["dimensions"][dn]["score"])
                for r in rows
                if r.get("dimensions", {}).get(dn, {}).get("score") is not None
            ]
            if mine is None or not scored:
                continue
            best = max(scored)
            if float(mine) >= best:
                lead.append(dn)
            elif best - float(mine) >= 2:
                trail.append(dn)
    d_obs = []
    if lead:
        d_obs.append(
            f"We hold the top score on {len(lead)} dimension{'s' if len(lead) != 1 else ''}: "
            + ", ".join(lead[:3])
            + "."
        )
    if trail:
        d_obs.append("We trail by two or more points on: " + ", ".join(trail[:3]) + ".")
    out["dimensions"] = {"observations": d_obs, "implications": []}
    out["versus"] = {"observations": [], "implications": []}

    o_obs = []
    if offers:
        with_welcome = [o for o in offers if o.get("welcome")]
        o_obs.append(
            f"{len(with_welcome)} of {len(offers)} brands lead with a stated welcome offer."
        )
    out["offers"] = {"observations": o_obs, "implications": []}
    out["exhibits"] = {"observations": [], "implications": []}
    out["coverage"] = {"observations": [], "implications": []}
    return out


_EVENT_RE = re.compile(
    r"\b(is closing|will close|closing on|closes on|shut(?:ting)? down|ceas(?:e|es|ing) "
    r"operations|exit(?:s|ed|ing)? (?:the )?(?:market|state)|leav(?:es|ing) (?:the )?"
    r"(?:market|state)|acquired by|acquisition of|merg(?:es|ed|ing) with|rebrand(?:s|ed|ing)?"
    r"(?: to| as)|now available in|launch(?:es|ed|ing)? in|cease[- ]and[- ]desist|"
    r"regulator|banned in|no longer (?:available|accept))\b",
    re.I,
)


def market_events(
    evidence: list[dict[str, Any]], *, limit: int = 4
) -> list[dict[str, Any]]:
    """Corporate / market events on record — closures, exits, acquisitions,
    rebrands, launches, regulatory notices — read straight from the claims.
    A baseline pack has no diff to surface them through, and a competitor
    closing (LuckyLand, 2026-08-16: "closing September 14, 2026" on the
    homepage) is the most material fact in the room; it must not sit in an
    appendix row while the market-moves slide says "no material change".
    Newest first, one per brand × claim."""
    out: list[dict[str, Any]] = []
    seen: set[tuple[str, str]] = set()
    for e in evidence:  # newest first
        claim = str(e.get("claim") or "").strip()
        if not claim or not _EVENT_RE.search(claim):
            continue
        brand = str(e.get("subject") or "").strip()
        # Two phrasings of one event ("…is closing on September 14, 2026" and
        # the same with a trailing clause) are one event: key on the opening.
        key = (brand, re.sub(r"[^a-z0-9]+", " ", claim.lower()).strip()[:44])
        if key in seen:
            continue
        seen.add(key)
        out.append(
            {
                "brand": brand,
                "claim": claim,
                "when": str(e.get("value_text") or ""),
                "observed_at": str(e.get("observed_at") or "")[:10],
                "url": str(e.get("source_url") or ""),
            }
        )
        if len(out) >= limit:
            break
    return out


def _events_banner(s: Any, events: list[dict[str, Any]], y: float) -> float:
    """A red-eyebrow strip under the headline: 'MARKET EVENT · <claim>'.
    Returns the y below the strip."""
    if not events:
        return y
    from pptx.util import Pt

    for ev in events[:2]:
        box = s.shapes.add_shape(1, _in(0.7), _in(y), _in(11.9), _in(0.36))
        box.fill.solid()
        box.fill.fore_color.rgb = _rgb(_CARD)
        box.line.color.rgb = _rgb(_ACCENT)
        box.line.width = Pt(0.75)
        box.shadow.inherit = False
        _text(s, 0.82, y + 0.05, 1.5, 0.28, "MARKET EVENT", size=8.5, bold=True, color=_ACCENT)
        line = f"{ev['brand']} – {ev['claim']}"
        if ev.get("observed_at"):
            line += f"  (observed {ev['observed_at']})"
        _text(s, 2.3, y + 0.05, 10.2, 0.28, _clean(line, 150), size=10, bold=True, color=_INK)
        y += 0.44
    return y


def factual_narrative(
    card: dict[str, Any],
    diff: dict[str, Any] | None,
    judged: list[dict[str, Any]],
    gaps: list[dict[str, Any]],
) -> dict[str, Any]:
    """The deck's words when no model is available: numbers only, no claims.

    Deliberately dull — a dull true summary beats a sharp invented one. Same
    shape as the model's narrative, so the renderer never branches; the
    model-only sections (exec zones, competitor profiles) stay empty and
    their slides are skipped rather than faked.
    """
    rows = card.get("rows", [])
    ranked = [r for r in rows if r.get("rank") is not None]
    us = next((r for r in rows if r.get("is_self")), None)
    leader = ranked[0] if ranked else None

    bullets: list[str] = []
    if leader:
        bullets.append(
            f"{leader['name']} leads the ranked field at "
            f"{_fmt(leader['overall']['normalized_pct'])} across {len(ranked)} "
            f"ranked brands ({len(rows)} tracked)."
        )
    else:
        bullets.append(
            f"{len(rows)} brands tracked; none yet scored on enough of the model to rank."
        )
    if us is not None:
        if us.get("rank") is not None:
            bullets.append(
                f"{us['name']} ranks #{us['rank']} at {_fmt(us['overall']['normalized_pct'])}."
            )
        elif us["overall"]["normalized_pct"] is not None:
            bullets.append(
                f"{us['name']} scores {_fmt(us['overall']['normalized_pct'])} "
                "but is provisional – not enough of the model measured to rank."
            )
        else:
            bullets.append(f"{us['name']} is not yet scored on any dimension.")
    if diff is None:
        bullets.append("First cycle: this pack sets the baseline; change appears next cycle.")
    else:
        n = int(diff.get("material_count", 0))
        bullets.append(
            f"{n} material change{'s' if n != 1 else ''} since the last snapshot."
            if n
            else "No material change since the last snapshot."
        )
    never = [g for g in gaps if g.get("status") == "never_observed"]
    stale = [g for g in gaps if g.get("status") == "stale"]
    bullets.append(
        f"{len(never)} brand × dimension pairs never observed; {len(stale)} overdue a refresh."
    )
    asks = [
        j
        for j in judged
        if str(j.get("decision_required", "none")).strip().lower() not in ("", "none")
    ]

    never_brands = sorted({g["subject"] for g in never})
    next_steps: list[str] = []
    if never_brands:
        more = " and others" if len(never_brands) > 3 else ""
        next_steps.append("Collect the unobserved brands – " + ", ".join(never_brands[:3]) + more)
    if stale:
        next_steps.append(f"Refresh the {len(stale)} pairs overdue against their cadence")
    if asks:
        plural = "s" if len(asks) != 1 else ""
        next_steps.append(f"Decide the {len(asks)} board ask{plural} on the decisions slide")
    if not next_steps:
        next_steps.append("Hold the cadence – re-observe on schedule and diff next cycle")

    # Executive-summary columns, computed: per top-weight dimension, who
    # holds the top score and where we stand. Dull but true.
    by_dim: list[dict[str, Any]] = []
    for d in sorted(card.get("dimensions", []), key=lambda d: -float(d.get("weight_pct") or 0))[:6]:
        dn = d["name"]
        scored = [
            (r, float(r["dimensions"][dn]["score"]))
            for r in rows
            if r.get("dimensions", {}).get(dn, {}).get("score") is not None
        ]
        obs: list[str] = []
        if scored:
            best = max(sc for _, sc in scored)
            names = [r["name"] for r, sc in scored if sc == best][:2]
            verb = "holds" if len(names) == 1 else "hold"
            obs.append(f"{' and '.join(names)} {verb} the top score ({best:g}/5).")
            if us is not None:
                mine = us.get("dimensions", {}).get(dn, {}).get("score")
                if mine is not None:
                    obs.append(f"{us['name']} scores {float(mine):g}/5.")
        else:
            obs.append("Not yet observed for any brand.")
        by_dim.append({"dimension": dn, "observations": obs})

    return {
        "headline": "",
        "bullets": bullets[:5],
        "exec": {"by_dimension": by_dim, "recommendation": "", "actions": []},
        "profiles": [],
        "titles": {},
        "commentary": {},
        "slides": _slides_facts(card),
        "next_steps": next_steps[:4],
        "source": "facts",
    }


# Back-compat alias for older callers.
factual_summary = factual_narrative


# ── slide builders ───────────────────────────────────────────────────


def _slide_title(
    prs: Any, *, title: str, market: str, period: str, basis: str, generated: str
) -> None:
    s = _blank(prs)
    _dark(s)
    _eyebrow(s, market or "Competitive intelligence", y=1.05)
    _text(
        s,
        0.7,
        1.55,
        11.9,
        1.9,
        _clean(title, 90),
        size=40,
        bold=True,
        color=_WHITE,
        line=1.06,
    )
    _text(s, 0.7, 3.55, 11.0, 0.5, _clean(period, 120), size=15, color=_MUTED)
    _rule(s, 4.35)
    _text(s, 0.7, 4.6, 11.9, 0.4, _clean(basis, 160), size=11, color=_MUTED)
    try:
        month = datetime.fromisoformat(generated).strftime("%B %Y")
    except Exception:
        month = datetime.now(UTC).strftime("%B %Y")
    _text(s, 0.7, 6.75, 5.0, 0.3, month, size=10, color=_MUTED)


def _slide_summary(
    prs: Any,
    card: dict[str, Any],
    narrative: dict[str, Any],
    page: int,
    deck_title: str,
    events: list[dict[str, Any]] | None = None,
) -> None:
    """The executive summary a steering committee reads first, in the shape
    of the reference decks: one column per battleground with numbered
    observations, then a recommendation strip — recommendation, where we
    stand, decisions. Nothing on this slide is a new fact; every line is a
    reading of the scorecard, the offers table or the material-change diff.

    Falls back to the three-zone findings/threats/watch layout when the
    narrative carries no per-dimension observations."""
    s = _blank(prs)
    _eyebrow(s, "Executive summary", y=0.42)
    headline = _clean(narrative.get("headline") or "", 120)
    y = 0.72
    if headline:
        _text(s, 0.7, y, 11.9, 0.9, headline, size=21, bold=True, line=1.08)
        y += 0.9
    else:
        y += 0.2
    y = _events_banner(s, events or [], y)
    _rule(s, y)
    y += 0.22

    exec_zone = narrative.get("exec") or {}
    by_dim = [
        d
        for d in (exec_zone.get("by_dimension") or [])
        if isinstance(d, dict) and str(d.get("dimension") or "").strip()
    ][:6]
    rows = card.get("rows", [])
    us_rows = [r for r in rows if r.get("is_self")]
    ranked = [r for r in rows if r.get("rank") is not None]
    leader = next((r for r in ranked if not r.get("is_self")), None)

    if by_dim:
        # ── dimension columns ──
        n = len(by_dim)
        gutter = 0.18
        col_w = (11.9 - gutter * (n - 1)) / n
        band_top = y
        band_bottom = 4.55
        counter = 0
        for ci, d in enumerate(by_dim):
            x = 0.7 + ci * (col_w + gutter)
            # column header
            hdr = s.shapes.add_shape(1, _in(x), _in(band_top), _in(col_w), _in(0.5))
            hdr.fill.solid()
            hdr.fill.fore_color.rgb = _rgb(_CARD)
            hdr.line.fill.background()
            hdr.shadow.inherit = False
            _text(
                s,
                x + 0.08,
                band_top + 0.06,
                col_w - 0.16,
                0.42,
                _clean(str(d.get("dimension")), 46),
                size=9.5,
                bold=True,
                color=_INK,
                line=1.0,
            )
            oy = band_top + 0.62
            for obs in [str(o) for o in (d.get("observations") or []) if str(o).strip()][:3]:
                if oy > band_bottom - 0.4:
                    break
                counter += 1
                _chip(
                    s,
                    x,
                    oy + 0.02,
                    counter,
                    color=_ACCENT if any(u["name"] in obs for u in us_rows) else _PEER,
                )
                text = _clean(obs, 130)
                lines = 1 + len(text) // max(18, int(col_w * 9.5))
                h = min(0.22 * lines + 0.1, band_bottom - oy)
                _text(s, x + 0.34, oy, col_w - 0.36, h, text, size=9, color=_BODY, line=1.05)
                oy += h + 0.1
        # ── recommendation strip ──
        y = band_bottom + 0.2
        _rule(s, y)
        y += 0.2
        strip_h = 6.5 - y
        boxes = [
            ("Recommendation", 0.7, 4.7),
            ("Where we stand", 5.55, 3.55),
            ("Decisions / next steps", 9.25, 3.35),
        ]
        rec = _clean(exec_zone.get("recommendation") or headline or "", 260)
        stand: list[str] = []
        for u in us_rows[:2]:
            if u.get("rank") is not None:
                line = f"{u['name']} – #{u['rank']} at {_fmt(u['overall']['normalized_pct'])}"
                if leader is not None:
                    gap = float(leader["overall"]["normalized_pct"]) - float(
                        u["overall"]["normalized_pct"]
                    )
                    line += f", {gap:.1f} behind {leader['name']}"
                stand.append(line)
            elif u["overall"]["normalized_pct"] is not None:
                stand.append(
                    f"{u['name']} – scores {_fmt(u['overall']['normalized_pct'])}, not yet ranked"
                )
            else:
                stand.append(f"{u['name']} – not yet scored")
        if leader is not None:
            stand.append(f"Leader: {leader['name']} at {_fmt(leader['overall']['normalized_pct'])}")
        if card.get("comparability_note"):
            stand.append("Ranks withheld until coverage is comparable")
        actions = [str(a) for a in (exec_zone.get("actions") or []) if str(a).strip()][:3] or [
            str(a) for a in (narrative.get("next_steps") or [])
        ][:3]
        for label, x, _w in boxes:
            _eyebrow(s, label, y=y, x=x, color=_ACCENT if label == "Recommendation" else _MUTED)
        if rec:
            _text(
                s, 0.7, y + 0.35, 4.7, strip_h - 0.4, rec, size=11, bold=True, color=_INK, line=1.12
            )
        _bullets(
            s,
            5.55,
            y + 0.35,
            3.55,
            strip_h - 0.4,
            stand or ["No brand marked as ours."],
            size=9.5,
            color=_BODY,
            gap_pt=5,
            cap=110,
            accent_bullet=False,
            max_items=4,
        )
        _bullets(
            s,
            9.25,
            y + 0.35,
            3.35,
            strip_h - 0.4,
            actions or ["Hold the cadence; diff next cycle."],
            size=9.5,
            color=_BODY,
            gap_pt=5,
            cap=110,
            accent_bullet=True,
            max_items=3,
        )
    else:
        findings = [str(x) for x in exec_zone.get("findings") or []]
        threats = [str(x) for x in exec_zone.get("threats") or []]
        watch = [str(x) for x in exec_zone.get("watch") or []]
        if findings or threats or watch:
            cols = [
                ("Key findings", _PEER, findings, 0.7, 3.9),
                ("Key threats", _ACCENT, threats, 4.95, 3.9),
                ("Watch next", _MUTED, watch, 9.2, 3.4),
            ]
            for label, label_color, items, x, w in cols:
                _eyebrow(s, label, y=y, x=x, color=label_color)
                _bullets(
                    s,
                    x,
                    y + 0.38,
                    w,
                    6.35 - y,
                    items or ["Nothing this period."],
                    size=12,
                    gap_pt=9,
                    cap=140,
                    accent_bullet=False,
                    max_items=4,
                )
        else:
            _bullets(
                s,
                0.7,
                y,
                11.9,
                6.4 - y,
                [str(b) for b in narrative.get("bullets") or []],
                size=15,
                gap_pt=12,
                cap=140,
            )
    _judgement_note(s, str(narrative.get("source") or "facts"))
    _footer(s, deck_title, page)
    _notes(
        s,
        "Every line traces to the scorecard, the offers table, the material-change "
        "diff or the evidence register. Nothing on this slide is a new fact. Numbers "
        "run across the columns in reading order.",
    )


def _in(v: float) -> Any:
    from pptx.util import Inches

    return Inches(v)


def _slide_glance(
    prs: Any,
    card: dict[str, Any],
    evidence_count: int,
    gaps: list[dict[str, Any]],
    commentary: str,
    page: int,
    deck_title: str,
) -> None:
    s = _blank(prs)
    _header(s, "The market at a glance", "The numbers under everything else")
    rows = card.get("rows", [])
    dims = card.get("dimensions", [])
    ranked = [r for r in rows if r.get("rank") is not None]
    leader = ranked[0] if ranked else None
    pairs = len(rows) * len(dims)
    observed = sum(
        1 for r in rows for d in r.get("dimensions", {}).values() if d.get("score") is not None
    )
    pct = (observed / pairs * 100.0) if pairs else 0.0

    scored = sorted(
        (r for r in rows if r["overall"]["normalized_pct"] is not None),
        key=lambda r: -float(r["overall"]["normalized_pct"]),
    )
    # Ranks withheld (comparability) or nothing ranked yet: the scores are
    # still the numbers under everything else — show the highest and say it
    # is unranked, rather than a dash and "0 / 14 ranked" (2026-08-16 pack).
    if leader is not None:
        first = (
            _fmt(leader["overall"]["normalized_pct"]),
            f"{leader['name']} – ranked leader",
        )
        second = (f"{len(ranked)} / {len(rows)}", "brands ranked / tracked")
    elif scored:
        first = (
            _fmt(scored[0]["overall"]["normalized_pct"]),
            f"{scored[0]['name']} – highest score, ranks withheld"
            if card.get("comparability_note")
            else f"{scored[0]['name']} – highest score, provisional",
        )
        second = (f"{len(scored)} / {len(rows)}", "brands scored / tracked – none ranked yet")
    else:
        first = ("—", "no brand scored yet")
        second = (f"0 / {len(rows)}", "brands scored / tracked")
    tiles = [
        first,
        second,
        (f"{pct:.0f}%", f"of the model observed – {observed} of {pairs} pairs"),
        (f"{evidence_count:,}", "observed facts, each traceable to a source"),
    ]
    x = 0.7
    w = 12.0 / len(tiles)
    for value, label in tiles:
        _text(s, x, 2.6, w - 0.35, 1.0, str(value)[:14], size=34, bold=True)
        _text(s, x, 3.62, w - 0.35, 0.75, _clean(label, 70), size=10.5, color=_MUTED)
        x += w
    if commentary:
        _text(
            s,
            0.7,
            5.0,
            11.9,
            0.5,
            _clean(commentary, 170),
            size=12.5,
            italic=True,
            color=_BODY,
        )
    _footer(s, deck_title, page)


def _slide_standings(
    prs: Any,
    card: dict[str, Any],
    narrative: dict[str, Any],
    page: int,
    deck_title: str,
) -> None:
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
    from pptx.util import Inches, Pt

    s = _blank(prs)
    title = (narrative.get("titles") or {}).get("standings") or "Where the market stands"
    top = _header(
        s,
        "Standings",
        title,
        (narrative.get("commentary") or {}).get("standings", ""),
    )

    rows = card.get("rows", [])
    ranked = [r for r in rows if r.get("rank") is not None]
    provisional = [
        r for r in rows if r.get("provisional") and r["overall"]["normalized_pct"] is not None
    ]
    unscored = [r for r in rows if r["overall"]["normalized_pct"] is None]

    if ranked:
        shown = ranked[:12]
        cd = CategoryChartData()
        cd.categories = [
            f"{r['name']}{'  (us)' if r.get('is_self') else ''}" for r in reversed(shown)
        ]
        cd.add_series(
            "Overall (normalized %)",
            [round(float(r["overall"]["normalized_pct"]), 1) for r in reversed(shown)],
        )
        gf = s.shapes.add_chart(
            XL_CHART_TYPE.BAR_CLUSTERED,
            Inches(0.7),
            Inches(top),
            Inches(8.1),
            Inches(6.6 - top),
            cd,
        )
        ch = gf.chart
        ch.has_legend = False
        ch.has_title = False
        va = ch.value_axis
        va.minimum_scale = 0
        va.maximum_scale = 100
        va.has_major_gridlines = False
        va.tick_labels.font.size = Pt(9)
        va.tick_labels.font.color.rgb = _rgb(_MUTED)
        ca = ch.category_axis
        ca.tick_labels.font.size = Pt(11)
        ca.tick_labels.font.color.rgb = _rgb(_INK)
        plot = ch.plots[0]
        plot.gap_width = 55
        plot.has_data_labels = True
        dl = plot.data_labels
        dl.font.size = Pt(10)
        dl.font.bold = True
        dl.font.color.rgb = _rgb(_INK)
        dl.number_format = "0.0"
        dl.number_format_is_linked = False
        dl.position = XL_LABEL_POSITION.OUTSIDE_END
        ser = plot.series[0]
        ser.format.fill.solid()
        ser.format.fill.fore_color.rgb = _rgb(_PEER)
        for idx, r in enumerate(reversed(shown)):
            if r.get("is_self"):
                pt = ser.points[idx]
                pt.format.fill.solid()
                pt.format.fill.fore_color.rgb = _rgb(_ACCENT)
        if len(ranked) > len(shown):
            _text(
                s,
                0.7,
                6.62,
                8.1,
                0.3,
                f"Top {len(shown)} of {len(ranked)} ranked brands shown.",
                size=9,
                color=_MUTED,
            )
    elif card.get("comparability_note") and provisional:
        # Ranks withheld, scores real: say so in one line, then chart the
        # scores unranked — the room still needs to see the numbers.
        _text(
            s,
            0.7,
            top,
            8.1,
            0.55,
            "Not yet a league table – " + _clean(card["comparability_note"], 150) + ".",
            size=9.5,
            italic=True,
            color=_ACCENT,
        )
        shown = sorted(provisional, key=lambda r: -float(r["overall"]["normalized_pct"]))[:14]
        cd = CategoryChartData()
        cd.categories = [
            f"{r['name']}{'  (us)' if r.get('is_self') else ''}" for r in reversed(shown)
        ]
        cd.add_series(
            "Score (unranked)",
            [round(float(r["overall"]["normalized_pct"]), 1) for r in reversed(shown)],
        )
        gf = s.shapes.add_chart(
            XL_CHART_TYPE.BAR_CLUSTERED,
            Inches(0.7),
            Inches(top + 0.55),
            Inches(8.1),
            Inches(6.6 - top - 0.55),
            cd,
        )
        ch = gf.chart
        ch.has_legend = False
        ch.has_title = False
        va = ch.value_axis
        va.minimum_scale = 0
        va.maximum_scale = 100
        va.has_major_gridlines = False
        va.tick_labels.font.size = Pt(9)
        va.tick_labels.font.color.rgb = _rgb(_MUTED)
        ca = ch.category_axis
        ca.tick_labels.font.size = Pt(10)
        ca.tick_labels.font.color.rgb = _rgb(_INK)
        plot = ch.plots[0]
        plot.gap_width = 55
        plot.has_data_labels = True
        plot.data_labels.font.size = Pt(9)
        plot.data_labels.font.color.rgb = _rgb(_BODY)
        plot.data_labels.number_format = "0.0"
        plot.data_labels.number_format_is_linked = False
        plot.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
        ser = plot.series[0]
        for i, r in enumerate(reversed(shown)):
            pt = ser.points[i]
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = _rgb(_ACCENT if r.get("is_self") else _PEER)
        provisional = []  # charted; do not repeat them in the side notes
    else:
        _text(
            s,
            0.7,
            top + 0.2,
            8.1,
            1.0,
            "No brand has yet been scored on enough of the model to hold a rank.",
            size=15,
            color=_BODY,
        )

    x = 9.1
    panel = (narrative.get("slides") or {}).get("standings") or {}
    y = _sidebar(
        s,
        panel.get("observations") or [],
        panel.get("implications") or [],
        top=top,
        x=x,
        w=3.5,
        bottom=(
            min(6.55, top + 2.3)
            if (provisional and unscored)
            else min(6.55, top + 3.0)
            if (provisional or unscored)
            else 6.55
        ),
    )
    y += 0.25 if y > top else 0.0
    if y > 6.3:
        y = 6.72  # sidebar used the column; the note shares the footnote line
    _text(
        s,
        x,
        y,
        3.5,
        0.3 if y >= 6.7 else 0.5,
        "Overall score, normalized to the weight actually measured. Amber – our brand.",
        size=8 if y >= 6.7 else 8.5,
        color=_MUTED,
    )
    y += 0.55
    if provisional and y < 6.0:
        _eyebrow(s, "Scored, not ranked †", y=y, x=x)
        y += 0.32
        lines = [
            f"{r['name']}: {_fmt(r['overall']['normalized_pct'])} – "
            f"{r.get('provisional_reason', 'thin evidence')}"
            for r in provisional[: (2 if unscored else 4)]
        ]
        block_h = min(0.5 * len(lines) + 0.1, 6.55 - y - (0.7 if unscored else 0))
        _bullets(
            s,
            x,
            y,
            3.5,
            max(0.3, block_h),
            lines,
            size=9,
            color=_MUTED,
            gap_pt=5,
            cap=120,
            accent_bullet=False,
        )
        y += max(0.3, block_h) + 0.15
    if unscored and y < 6.4:
        _eyebrow(s, "Not yet observed", y=y, x=x)
        names = ", ".join(r["name"] for r in unscored[:8])
        if len(unscored) > 8:
            names += f" +{len(unscored) - 8} more"
        _text(s, x, y + 0.32, 3.5, min(1.4, 6.7 - y), names, size=9.5, color=_MUTED)
    _footer(s, deck_title, page)
    _notes(
        s,
        "Provisional brands are omitted from the chart on purpose: a bar is "
        "trusted without its footnote. Their scores are real; the standing is "
        "not yet earned. Unscored brands are never drawn as 0.",
    )


def _slide_versus(
    prs: Any,
    card: dict[str, Any],
    narrative: dict[str, Any],
    page: int,
    deck_title: str,
) -> None:
    s = _blank(prs)
    rows = card.get("rows", [])
    us = next((r for r in rows if r.get("is_self")), None)
    leader = next(
        (r for r in rows if r.get("rank") is not None and not r.get("is_self")),
        None,
    )
    if leader is None:
        # Ranks withheld (field not comparable) — compare against the
        # highest-scoring peer; dimension scores are still real.
        peers = [
            r for r in rows if not r.get("is_self") and r["overall"]["normalized_pct"] is not None
        ]
        leader = max(peers, key=lambda r: float(r["overall"]["normalized_pct"])) if peers else None
    title = (narrative.get("titles") or {}).get("versus") or "Where we win, where we lose"
    top = _header(
        s,
        "Versus the leader",
        title,
        (narrative.get("commentary") or {}).get("versus", ""),
    )

    if us is None:
        _text(
            s,
            0.7,
            top + 0.2,
            11.9,
            1.0,
            "No brand is marked as ours, so there is nothing to compare. Mark "
            "one with watch_subject is_self=true and this slide fills in.",
            size=14,
            color=_BODY,
        )
        _footer(s, deck_title, page)
        return
    if leader is None:
        _text(
            s,
            0.7,
            top + 0.2,
            11.9,
            1.0,
            f"{us['name']} is the only ranked brand so far – no peer holds a "
            "rank yet to compare against.",
            size=14,
            color=_BODY,
        )
        _footer(s, deck_title, page)
        return

    _text(
        s,
        0.7,
        top,
        11.9,
        0.45,
        f"{us['name']}  {_fmt(us['overall']['normalized_pct'])}   vs   "
        f"{leader['name']} "
        f"({'#' + str(leader['rank']) if leader.get('rank') is not None else 'highest-scoring peer'})  "
        f"{_fmt(leader['overall']['normalized_pct'])}",
        size=15,
        bold=True,
    )
    top += 0.6

    weights = {d["name"]: d["weight_pct"] for d in card.get("dimensions", [])}
    ahead: list[tuple[float, str]] = []
    behind: list[tuple[float, str]] = []
    unmeasured: list[str] = []
    for dname, ours in us.get("dimensions", {}).items():
        theirs = leader.get("dimensions", {}).get(dname, {})
        a, b = ours.get("score"), theirs.get("score")
        if a is None or b is None:
            unmeasured.append(dname)
            continue
        delta = float(a) - float(b)
        entry = f"{dname}  ({a:g} vs {b:g}, weight {weights.get(dname, 0):g}%)"
        if delta > 0:
            ahead.append((delta, entry))
        elif delta < 0:
            behind.append((-delta, entry))
    ahead.sort(key=lambda t: -t[0])
    behind.sort(key=lambda t: -t[0])

    cols = [
        (
            "Ahead",
            _ACCENT,
            [entry for _, entry in ahead[:5]] or ["Nowhere yet, on measured dimensions."],
            0.7,
            3.9,
            _BODY,
        ),
        (
            "Behind",
            _PEER,
            [entry for _, entry in behind[:5]] or ["Nowhere, on measured dimensions."],
            4.85,
            3.9,
            _BODY,
        ),
    ]
    for label, label_color, items, x, w, body_color in cols:
        _eyebrow(s, label, y=top, x=x, color=label_color)
        _bullets(
            s,
            x,
            top + 0.35,
            w,
            6.2 - top,
            items,
            size=10.5,
            color=body_color,
            gap_pt=6,
            cap=110,
            accent_bullet=False,
            max_items=7,
        )
    panel = (narrative.get("slides") or {}).get("versus") or {}
    _sidebar(
        s,
        panel.get("observations") or [],
        panel.get("implications") or [],
        top=top,
        x=9.1,
        w=3.5,
    )
    note = (
        "A dimension unscored on either side is listed as not comparable – never counted as a loss."
    )
    if unmeasured:
        note = (
            "Not comparable yet: "
            + ", ".join(unmeasured[:5])
            + (f" +{len(unmeasured) - 5}" if len(unmeasured) > 5 else "")
            + ". "
            + note
        )
    _text(s, 0.7, 6.62, 11.9, 0.45, _clean(note, 220), size=8.5, color=_MUTED)
    _footer(s, deck_title, page)


def _slide_dimension_leaders(
    prs: Any,
    card: dict[str, Any],
    narrative: dict[str, Any],
    page: int,
    deck_title: str,
) -> None:
    """Who leads each battleground: per weighted dimension, the top observed
    score, our score, and the gap — the per-area breakout a steering
    committee expects after the overall standings."""
    from pptx.util import Inches, Pt

    s = _blank(prs)
    title = (narrative.get("titles") or {}).get("dimensions") or "Who leads each dimension"
    top = _header(
        s,
        "The battlegrounds",
        title,
        (narrative.get("commentary") or {}).get("dimensions", ""),
    )
    rows = card.get("rows", [])
    dims = card.get("dimensions", [])[:12]
    us = next((r for r in rows if r.get("is_self")), None)
    if not dims or not rows:
        _text(
            s,
            0.7,
            top + 0.2,
            11.9,
            1.0,
            "Nothing scored yet.",
            size=15,
            color=_BODY,
        )
        _footer(s, deck_title, page)
        return

    def brand_label(r: dict[str, Any]) -> str:
        name = str(r["name"])
        if r.get("is_self"):
            name += " (us)"
        if r.get("provisional") and r["overall"]["normalized_pct"] is not None:
            name += " †"
        return name

    cols = ["Dimension", "Market leader", "Theirs", "Us", "Gap"]
    widths = [2.75, 3.05, 0.65, 0.65, 1.1]
    shape = s.shapes.add_table(
        len(dims) + 1,
        len(cols),
        Inches(0.7),
        Inches(top),
        Inches(sum(widths)),
        Inches(min(0.34 * (len(dims) + 1), 6.55 - top)),
    )
    tbl = shape.table
    for ci, w in enumerate(widths):
        tbl.columns[ci].width = Inches(w)

    def cell_write(
        r: int,
        c: int,
        text: str,
        *,
        bg: str,
        fg: str = _INK,
        bold: bool = False,
        size: float = 9.5,
    ) -> None:
        cell = tbl.cell(r, c)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = _rgb(bg)
        cell.margin_left = cell.margin_right = Inches(0.06)
        cell.margin_top = cell.margin_bottom = Inches(0.02)
        for p_ in cell.text_frame.paragraphs:
            for run in p_.runs:
                run.font.size = Pt(size)
                run.font.bold = bold
                run.font.color.rgb = _rgb(fg)

    for ci, name in enumerate(cols):
        cell_write(0, ci, name, bg=_INK, fg=_WHITE, bold=True, size=10)
    for ri, d in enumerate(dims, start=1):
        dname = d["name"]
        bg = _WHITE if ri % 2 else _CARD
        scored = [
            (r, float(r["dimensions"][dname]["score"]))
            for r in rows
            if r.get("dimensions", {}).get(dname, {}).get("score") is not None
        ]
        cell_write(
            ri,
            0,
            f"{dname}  ·  {d.get('weight_pct', 0):g}%",
            bg=bg,
            bold=True,
            size=9,
        )
        if not scored:
            cell_write(ri, 1, "not yet observed", bg=bg, fg=_MUTED)
            cell_write(ri, 2, "", bg=_GAP_BG)
            cell_write(ri, 3, "", bg=_GAP_BG)
            cell_write(ri, 4, "", bg=_GAP_BG)
            continue
        best = max(sc for _, sc in scored)
        leaders = [r for r, sc in scored if sc == best]
        names = ", ".join(brand_label(r) for r in leaders[:2])
        if len(leaders) > 2:
            names += f" +{len(leaders) - 2}"
        we_lead = us is not None and any(r.get("is_self") for r in leaders)
        cell_write(
            ri,
            1,
            names,
            bg=bg,
            fg=_ACCENT if we_lead else _INK,
            bold=we_lead,
            size=9,
        )
        cell_write(ri, 2, f"{best:g}", bg=bg)
        ours = us.get("dimensions", {}).get(dname, {}).get("score") if us is not None else None
        if us is None:
            cell_write(ri, 3, "—", bg=bg, fg=_MUTED)
            cell_write(ri, 4, "—", bg=bg, fg=_MUTED)
        elif ours is None:
            cell_write(ri, 3, "", bg=_GAP_BG)
            cell_write(ri, 4, "not measured", bg=bg, fg=_MUTED, size=8.5)
        elif we_lead:
            cell_write(ri, 3, f"{float(ours):g}", bg=bg, bold=True)
            others = [sc for r, sc in scored if not r.get("is_self")]
            margin = float(ours) - max(others) if others else None
            label = "we lead" if margin is None or margin > 0 else "we co-lead"
            cell_write(ri, 4, label, bg=bg, fg=_ACCENT, bold=True, size=8.5)
        else:
            delta = float(ours) - best
            cell_write(ri, 3, f"{float(ours):g}", bg=bg)
            cell_write(ri, 4, f"▼ {abs(delta):g} behind", bg=bg, size=8.5)
    panel = (narrative.get("slides") or {}).get("dimensions") or {}
    _sidebar(
        s,
        panel.get("observations") or [],
        panel.get("implications") or [],
        top=top,
        x=9.1,
        w=3.5,
    )
    if len(card.get("dimensions", [])) > len(dims):
        _text(
            s,
            0.7,
            6.6,
            8.2,
            0.3,
            f"First {len(dims)} of {len(card.get('dimensions', []))} dimensions "
            "shown — the rest are in the workbook.",
            size=9,
            color=_MUTED,
        )
    _text(
        s,
        0.7,
        6.85,
        11.9,
        0.25,
        "Top observed score per dimension. † overall standing provisional. "
        "A blank cell is unmeasured – never counted as a loss.",
        size=9,
        color=_MUTED,
    )
    _footer(s, deck_title, page)
    _notes(
        s,
        "Per-dimension leaders use dimension scores directly; a brand whose "
        "OVERALL is provisional can still hold a real top score on one "
        "dimension it was actually measured on — the dagger carries that "
        "context onto the slide.",
    )


def _slide_profile(
    prs: Any,
    row: dict[str, Any],
    profile: dict[str, Any],
    weights: dict[str, float],
    exhibits: list[dict[str, str]],
    page: int,
    deck_title: str,
) -> None:
    """One competitor, one slide: their scores and storefront on the left,
    what they are doing and what it means for us on the right — the
    steering-committee deep-dive pattern."""
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt

    s = _blank(prs)
    name = row["name"]
    rank = row.get("rank")
    overall = row["overall"]["normalized_pct"]
    standing = f"#{rank}" if rank is not None else "provisional †"
    eyebrow = f"Competitor deep dive · {standing} · overall {_fmt(overall)}"
    title = profile.get("title") or f"{name}"
    top = _header(s, eyebrow, f"{name} – {title}" if title != name else name)

    # Left: the brand's strongest-weighted scored dimensions as 1-5 bars.
    scored = [
        (dname, float(d["score"]))
        for dname, d in row.get("dimensions", {}).items()
        if d.get("score") is not None
    ]
    scored.sort(key=lambda t: -weights.get(t[0], 0.0))
    shown = scored[:5]
    not_observed = [
        dname for dname, d in row.get("dimensions", {}).items() if d.get("score") is None
    ]
    bar_x, bar_w = 0.7, 4.4
    y = top + 0.1
    for dname, score in shown:
        _text(s, bar_x, y, bar_w, 0.26, _clean(dname, 44), size=9.5, color=_BODY)
        track = s.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(bar_x),
            Inches(y + 0.26),
            Inches(bar_w),
            Pt(6),
        )
        track.fill.solid()
        track.fill.fore_color.rgb = _rgb(_GAP_BG)
        track.line.fill.background()
        track.shadow.inherit = False
        fill = s.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(bar_x),
            Inches(y + 0.26),
            Inches(bar_w * max(0.0, min(score, 5.0)) / 5.0),
            Pt(6),
        )
        fill.fill.solid()
        fill.fill.fore_color.rgb = _rgb(_ACCENT if row.get("is_self") else _PEER)
        fill.line.fill.background()
        fill.shadow.inherit = False
        _text(
            s,
            bar_x + bar_w + 0.08,
            y + 0.1,
            0.5,
            0.26,
            f"{score:g}",
            size=9.5,
            bold=True,
        )
        y += 0.52
    if not shown:
        _text(
            s,
            bar_x,
            y,
            bar_w,
            0.8,
            "No dimension scored yet – observations pending.",
            size=11,
            color=_MUTED,
        )
        y += 0.9
    extras: list[str] = []
    if len(scored) > len(shown):
        extras.append(f"+{len(scored) - len(shown)} more scored – see appendix")
    if not_observed:
        more = f" +{len(not_observed) - 3} more" if len(not_observed) > 3 else ""
        extras.append("not observed: " + ", ".join(not_observed[:3]) + more)
    if extras:
        _text(
            s,
            bar_x,
            y + 0.02,
            bar_w + 0.6,
            0.45,
            " · ".join(extras),
            size=8.5,
            color=_MUTED,
        )
        y += 0.45

    # Storefront capture under the bars — fixed slot so it always fits.
    if exhibits:
        shot = exhibits[0]
        shot_y = max(y + 0.15, 4.75)
        max_h = 6.6 - shot_y
        if max_h > 0.7:
            pic = _picture(s, shot["path"], bar_x, shot_y, 2.9, max_h)
            if pic is not None:
                cap = "Storefront capture"
                if shot.get("observed_at"):
                    cap += f" · {str(shot['observed_at'])[:10]}"
                _text(
                    s,
                    bar_x + 3.0,
                    shot_y + 0.05,
                    2.2,
                    0.8,
                    cap,
                    size=8.5,
                    color=_MUTED,
                )

    # Right: observations → implications, model-written from filed facts.
    x = 6.1
    _eyebrow(s, "Observations", y=top + 0.1, x=x, color=_PEER)
    _bullets(
        s,
        x,
        top + 0.45,
        6.4,
        2.9,
        [str(o) for o in profile.get("observations") or []]
        or ["No narrative available – see the workbook for this brand's facts."],
        size=11.5,
        gap_pt=8,
        cap=150,
        accent_bullet=False,
        max_items=4,
    )
    imp_y = top + 3.3
    _eyebrow(s, "Implications for us", y=imp_y, x=x, color=_ACCENT)
    _bullets(
        s,
        x,
        imp_y + 0.35,
        6.4,
        6.55 - imp_y - 0.4,
        [str(i) for i in profile.get("implications") or []] or ["Not yet judged."],
        size=11.5,
        gap_pt=8,
        cap=150,
        accent_bullet=False,
        max_items=2,
    )
    _judgement_note(s, "model")
    _footer(s, deck_title, page)
    _notes(
        s,
        "Bars are computed scores (1-5) on the weighted dimensions; blank "
        "means not observed, never zero. Observations and implications are "
        "model-written from this brand's filed facts.",
    )


def _slide_offers(
    prs: Any,
    offers: list[dict[str, Any]],
    narrative: dict[str, Any],
    page: int,
    deck_title: str,
) -> None:
    """The offers on the table: per brand, the headline welcome offer and
    the ongoing promotion, verbatim from the evidence register — what the
    market is actually selling to a new visitor this cycle. Executives ask
    for this first; it is the most comparable fact in the whole pack."""
    from pptx.util import Inches, Pt

    s = _blank(prs)
    title = (narrative.get("titles") or {}).get("offers") or "The offers on the table"
    top = _header(
        s,
        "Offers and promotions",
        title,
        (narrative.get("commentary") or {}).get("offers", ""),
    )
    shown = offers[:14]
    cols = ["Brand", "Headline welcome offer", "Ongoing / daily proposition"]
    widths = [1.9, 3.35, 2.95]
    tbl_h = min(0.36 * (len(shown) + 1), 6.5 - top)
    shape = s.shapes.add_table(
        len(shown) + 1, len(cols), Inches(0.7), Inches(top), Inches(sum(widths)), Inches(tbl_h)
    )
    tbl = shape.table
    for ci, w in enumerate(widths):
        tbl.columns[ci].width = Inches(w)

    def cell_write(
        r: int, c: int, text: str, *, bg: str, fg: str = _INK, bold: bool = False, size: float = 8.5
    ) -> None:
        cell = tbl.cell(r, c)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = _rgb(bg)
        cell.margin_left = cell.margin_right = Inches(0.05)
        cell.margin_top = cell.margin_bottom = Inches(0.02)
        for p_ in cell.text_frame.paragraphs:
            for run in p_.runs:
                run.font.size = Pt(size)
                run.font.bold = bold
                run.font.color.rgb = _rgb(fg)

    for ci, name in enumerate(cols):
        cell_write(0, ci, name, bg=_INK, fg=_WHITE, bold=True, size=9.5)
    for ri, o in enumerate(shown, start=1):
        is_self = bool(o.get("is_self"))
        bg = _SELF_ROW if is_self else (_WHITE if ri % 2 else _CARD)
        name = str(o.get("brand") or "") + ("  (us)" if is_self else "")
        cell_write(ri, 0, name, bg=bg, fg=_ACCENT if is_self else _INK, bold=True, size=8.5)
        cell_write(
            ri,
            1,
            _clean(o.get("welcome") or "not stated on the pages read", 95),
            bg=bg,
            fg=_INK if o.get("welcome") else _MUTED,
        )
        cell_write(
            ri,
            2,
            _clean(o.get("ongoing") or "—", 95),
            bg=bg,
            fg=_INK if o.get("ongoing") else _MUTED,
        )

    panel = (narrative.get("slides") or {}).get("offers") or {}
    _sidebar(
        s, panel.get("observations") or [], panel.get("implications") or [], top=top, x=9.1, w=3.5
    )
    _text(
        s,
        0.7,
        6.72,
        11.9,
        0.3,
        "Offer text as published on each brand's own pages at capture; amounts and "
        'conditions verified against the page excerpt. "Not stated" means the pages '
        "read carried no welcome offer, not that none exists.",
        size=8.5,
        color=_MUTED,
    )
    _footer(s, deck_title, page)
    _notes(
        s,
        "Every cell is a claim from the evidence register with a source URL and an "
        "excerpt; nothing here is paraphrased from memory.",
    )


def _slide_exhibits(
    prs: Any,
    batch: list[tuple[str, dict[str, str], bool]],
    idx: int,
    total: int,
    page: int,
    deck_title: str,
    *,
    eyebrow: str = "Exhibits",
    title: str = "The storefronts as a visitor sees them",
    captions: dict[str, str] | None = None,
    commentary: str = "",
) -> None:
    """Two captures per slide, photographed by the browser — the market as a
    visitor sees it, filed beside the claims it supports. ``captions`` adds
    a one-line reading under a brand's exhibit (e.g. the offer it shows)."""
    s = _blank(prs)
    suffix = f" ({idx}/{total})" if total > 1 else ""
    top = _header(s, eyebrow, f"{title}{suffix}", commentary)
    slots = [(0.7, 5.9), (6.85, 5.9)]
    max_h = 3.0 if captions else 3.4
    for (x, w), (brand, shot, is_self) in zip(slots, batch, strict=False):
        pic = _picture(s, shot["path"], x, top + 0.15, w, max_h)
        label = f"{brand}{'  (us)' if is_self else ''}"
        pic_h = (pic.height / 914400.0) if pic is not None else max_h
        cap_y = top + 0.15 + min(pic_h, max_h) + 0.12
        _text(s, x, cap_y, w, 0.3, label, size=12, bold=True, color=_ACCENT if is_self else _INK)
        line_y = cap_y + 0.3
        cap = (captions or {}).get(brand, "")
        if cap:
            _text(s, x, line_y, w, 0.45, _clean(cap, 140), size=9.5, color=_BODY)
            line_y += 0.42
        detail = []
        if shot.get("url"):
            detail.append(_clean(shot["url"], 70))
        if shot.get("observed_at"):
            detail.append(f"captured {str(shot['observed_at'])[:10]}")
        if detail:
            _text(s, x, line_y, w, 0.3, " · ".join(detail), size=8.5, color=_MUTED)
        if pic is None:
            _text(s, x, top + 1.5, w, 0.5, "exhibit unavailable", size=10, color=_MUTED)
    _text(
        s,
        0.7,
        6.72,
        11.9,
        0.3,
        "Captured by the browser through a state-verified network exit, cookie "
        "consent dismissed; each capture is filed in the evidence register beside its claims.",
        size=9,
        color=_MUTED,
    )
    _footer(s, deck_title, page)


def _slide_changes(
    prs: Any,
    diff: dict[str, Any] | None,
    narrative: dict[str, Any],
    page: int,
    deck_title: str,
    events: list[dict[str, Any]] | None = None,
) -> None:
    s = _blank(prs)
    commentary = (narrative.get("commentary") or {}).get("changes", "")
    events = events or []
    if diff is None:
        n = len(events)
        title = (
            f"Baseline established – {n} market event{'s' if n != 1 else ''} on record"
            if n
            else "Baseline established"
        )
        top = _header(s, "Market moves", title, commentary)
        _text(
            s,
            0.7,
            top,
            11.9,
            0.8,
            "First reporting cycle. There is no prior snapshot to compare "
            "against, so this pack sets the baseline; the next cycle shows "
            "what moved."
            + (
                " The events below were read from the brands' own pages during "
                "this cycle and need no prior snapshot."
                if events
                else ""
            ),
            size=13,
            color=_BODY,
        )
        if events:
            _eyebrow(s, "Market events on record", y=top + 0.95, color=_ACCENT)
            _bullets(
                s,
                0.7,
                top + 1.3,
                11.9,
                6.4 - (top + 1.3),
                [
                    f"{ev['brand']} – {ev['claim']}"
                    + (f" (observed {ev['observed_at']})" if ev.get("observed_at") else "")
                    for ev in events
                ],
                size=12,
                color=_INK,
                gap_pt=8,
                cap=200,
                accent_bullet=True,
                max_items=4,
            )
        _footer(s, deck_title, page)
        return
    n = int(diff.get("material_count", 0))
    title = (narrative.get("titles") or {}).get(
        "changes"
    ) or f"{n} material move{'s' if n != 1 else ''} this period"
    top = _header(s, "Market moves", title, commentary)
    _text(
        s,
        0.7,
        top - 0.05,
        11.9,
        0.3,
        f"{str(diff.get('from_generated_at', '?'))[:10]} → "
        f"{str(diff.get('to_generated_at', '?'))[:10]}",
        size=10,
        color=_MUTED,
    )
    top += 0.3
    if n == 0:
        th = diff.get("thresholds", {})
        _text(
            s,
            0.7,
            top + 0.2,
            11.9,
            1.2,
            "Nothing moved above the materiality threshold "
            f"(score move ≥ {th.get('min_score_delta', 1.0)}, coverage shift ≥ "
            f"{th.get('min_coverage_delta', 0)}%). Sub-threshold wobble is "
            "suppressed on purpose.",
            size=15,
            color=_BODY,
        )
        _footer(s, deck_title, page)
        return
    lines: list[str] = []
    for c in diff.get("changed", []):
        for item in c.get("items", []):
            lines.append(f"{c['subject']} – {item.get('detail', '')}")
    if diff.get("added_subjects"):
        lines.append("New in the analysis: " + ", ".join(diff["added_subjects"]))
    if diff.get("removed_subjects"):
        lines.append("Removed: " + ", ".join(diff["removed_subjects"]))
    shown = lines[:11]
    half = (len(shown) + 1) // 2
    _bullets(
        s,
        0.7,
        top,
        5.9,
        6.5 - top,
        shown[:half],
        size=11.5,
        gap_pt=7,
        cap=120,
        accent_bullet=False,
        max_items=6,
    )
    _bullets(
        s,
        6.85,
        top,
        5.8,
        6.5 - top,
        shown[half:],
        size=11.5,
        gap_pt=7,
        cap=120,
        accent_bullet=False,
        max_items=6,
    )
    if len(lines) > len(shown):
        _text(
            s,
            0.7,
            6.55,
            11.9,
            0.3,
            f"+{len(lines) - len(shown)} more in the board report.",
            size=9,
            color=_MUTED,
        )
    _footer(s, deck_title, page)


def _slide_implications(
    prs: Any,
    judged: list[dict[str, Any]],
    material_count: int,
    source: str,
    page: int,
    deck_title: str,
) -> int:
    from pptx.util import Inches, Pt

    if not judged:
        s = _blank(prs)
        _header(s, "Judgement", "Implications and recommendations")
        msg = (
            "Implications were not generated – no model was available. The "
            "material changes are the factual record; judgement still needs "
            "to be applied."
            if material_count
            else "No material change this period, so there is nothing new to act on."
        )
        _text(s, 0.7, 2.4, 11.9, 1.2, msg, size=15, color=_BODY)
        _footer(s, deck_title, page)
        return 1

    ordered = sorted(
        judged,
        key=lambda j: (
            _CLASS_ORDER.index(j.get("classification", "monitor"))
            if j.get("classification") in _CLASS_ORDER
            else len(_CLASS_ORDER)
        ),
    )
    per = 5
    batches = [ordered[i : i + per] for i in range(0, len(ordered), per)]
    for pi, chunk in enumerate(batches):
        s = _blank(prs)
        suffix = f" ({pi + 1}/{len(batches)})" if len(batches) > 1 else ""
        top = _header(s, "Judgement", f"Implications and recommendations{suffix}")
        cols = ["Subject", "Change", "Implication", "Recommendation", "Class"]
        widths = [1.7, 2.7, 3.3, 3.3, 1.6]
        shape = s.shapes.add_table(
            len(chunk) + 1,
            len(cols),
            Inches(0.7),
            Inches(top),
            Inches(sum(widths)),
            Inches(0.42 * (len(chunk) + 1)),
        )
        tbl = shape.table
        for ci, w in enumerate(widths):
            tbl.columns[ci].width = Inches(w)
        for ci, name in enumerate(cols):
            cell = tbl.cell(0, ci)
            cell.text = name
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb(_INK)
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(10)
                    r.font.bold = True
                    r.font.color.rgb = _rgb(_WHITE)
        for ri, j in enumerate(chunk, start=1):
            vals = [
                _clean(j.get("subject", ""), 40),
                _clean(j.get("change", ""), 90),
                _clean(j.get("implication", ""), 110),
                _clean(j.get("recommendation", ""), 110),
                _CLASS_LABEL.get(str(j.get("classification", "monitor")), "Monitor"),
            ]
            for ci, v in enumerate(vals):
                cell = tbl.cell(ri, ci)
                cell.text = v
                cell.fill.solid()
                cell.fill.fore_color.rgb = _rgb(_WHITE if ri % 2 else _CARD)
                for p in cell.text_frame.paragraphs:
                    for r in p.runs:
                        r.font.size = Pt(9)
                        r.font.color.rgb = _rgb(_INK)
        _judgement_note(s, source)
        _footer(s, deck_title, page + pi)
        _notes(
            s,
            "Classification uses the no-regret test: would we still be pleased "
            "we did this if the transition plan, rankings or priorities "
            "changed next month?",
        )
    return len(batches)


def _slide_decisions(
    prs: Any,
    judged: list[dict[str, Any]],
    source: str,
    material_count: int,
    page: int,
    deck_title: str,
) -> None:
    s = _blank(prs)
    top = _header(s, "The ask", "Decisions required")
    asks = [
        f"{j.get('subject', '')}: {j.get('decision_required', '')}"
        for j in judged
        if str(j.get("decision_required", "none")).strip().lower() not in ("", "none")
    ]
    if asks:
        _bullets(s, 0.7, top + 0.2, 11.9, 6.4 - top, asks[:6], size=15, gap_pt=12, cap=140)
        _judgement_note(s, source)
    elif judged:
        _text(
            s,
            0.7,
            top + 0.3,
            11.9,
            1.0,
            "No board decision is required this period.",
            size=17,
            color=_BODY,
        )
        _judgement_note(s, source)
    elif material_count:
        _text(
            s,
            0.7,
            top + 0.3,
            11.9,
            1.4,
            f"Not yet evaluated – {material_count} material change"
            f"{'s' if material_count != 1 else ''} recorded, but no model was "
            "available to judge what they require of the board.",
            size=17,
            color=_BODY,
        )
    else:
        _text(
            s,
            0.7,
            top + 0.3,
            11.9,
            1.0,
            "No material change this period, so nothing new to decide.",
            size=17,
            color=_BODY,
        )
    _footer(s, deck_title, page)


def _slide_evidence(
    prs: Any,
    card: dict[str, Any],
    gaps: list[dict[str, Any]],
    evidence_count: int,
    narrative: dict[str, Any],
    page: int,
    deck_title: str,
) -> None:
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt

    s = _blank(prs)
    title = (narrative.get("titles") or {}).get("coverage") or "How much of this is measured"
    top = _header(
        s,
        "Appendix – evidence and confidence",
        title,
        (narrative.get("commentary") or {}).get("coverage", ""),
    )
    rows = card.get("rows", [])
    dims = card.get("dimensions", [])
    pairs = len(rows) * len(dims)
    observed = sum(
        1 for r in rows for d in r.get("dimensions", {}).values() if d.get("score") is not None
    )
    never = [g for g in gaps if g.get("status") == "never_observed"]
    stale = [g for g in gaps if g.get("status") == "stale"]
    pct = (observed / pairs * 100.0) if pairs else 0.0

    tiles = [
        (f"{pct:.0f}%", f"of the model observed – {observed} of {pairs} brand × dimension pairs"),
        (f"{evidence_count:,}", "evidence items on file, each quoting a source URL"),
        (f"{len(never)}", "pairs never observed – no score, no penalty"),
        (f"{len(stale)}", "overdue a refresh against their cadence"),
    ]
    x = 0.7
    for big, label in tiles:
        _text(s, x, top + 0.15, 2.85, 0.9, big, size=32, bold=True)
        _text(s, x, top + 1.05, 2.85, 0.8, _clean(label, 80), size=10, color=_MUTED)
        x += 3.05

    y = top + 2.15
    never_brands = sorted({g["subject"] for g in never})
    if never_brands:
        more = f" +{len(never_brands) - 6} more" if len(never_brands) > 6 else ""
        _text(
            s,
            0.7,
            y,
            11.9,
            0.35,
            "Not yet observed: " + ", ".join(never_brands[:6]) + more,
            size=11.5,
            color=_BODY,
        )
        y += 0.45
    if stale:
        ex = ", ".join(f"{g['subject']} / {g['dimension']}" for g in stale[:4])
        _text(s, 0.7, y, 11.9, 0.35, f"Overdue: {ex}", size=11.5, color=_BODY)
        y += 0.45

    card_y = max(y + 0.15, 5.5)
    box = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(card_y), Inches(11.9), Inches(0.85)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = _rgb(_CARD)
    box.line.color.rgb = _rgb(_HAIR)
    box.line.width = Pt(0.75)
    box.shadow.inherit = False
    bar = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(card_y), Inches(0.06), Inches(0.85)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = _rgb(_ACCENT)
    bar.line.fill.background()
    bar.shadow.inherit = False
    _text(
        s,
        0.95,
        card_y + 0.12,
        11.4,
        0.65,
        "Gaps are absences of evidence, not weaknesses. No brand is scored "
        "down for being opaque, and no score on any slide exists without a "
        "source behind it.",
        size=11.5,
        color=_INK,
    )
    _footer(s, deck_title, page)


def _slide_heatmap(prs: Any, card: dict[str, Any], page: int, deck_title: str) -> None:
    from pptx.util import Inches, Pt

    s = _blank(prs)
    top = _header(s, "Appendix", "Scores by dimension")
    rows = card.get("rows", [])[:14]
    dims = [d["name"] for d in card.get("dimensions", [])][:12]
    if not rows or not dims:
        _text(s, 0.7, top + 0.2, 11.9, 1.0, "Nothing scored yet.", size=15, color=_BODY)
        _footer(s, deck_title, page)
        return
    ncols = 2 + len(dims)
    total_w = 11.9
    first_w = 2.0
    dim_w = (total_w - first_w - 0.9) / len(dims)
    shape = s.shapes.add_table(
        len(rows) + 1,
        ncols,
        Inches(0.7),
        Inches(top),
        Inches(total_w),
        Inches(min(0.32 * (len(rows) + 1), 6.6 - top)),
    )
    tbl = shape.table
    tbl.columns[0].width = Inches(first_w)
    tbl.columns[1].width = Inches(0.9)
    for ci in range(len(dims)):
        tbl.columns[2 + ci].width = Inches(dim_w)

    def cell_write(
        r: int, c: int, text: str, *, bg: str, fg: str = _INK, bold: bool = False, size: int = 8
    ) -> None:
        cell = tbl.cell(r, c)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = _rgb(bg)
        cell.margin_left = cell.margin_right = Inches(0.03)
        cell.margin_top = cell.margin_bottom = Inches(0.01)
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                run.font.size = Pt(size)
                run.font.bold = bold
                run.font.color.rgb = _rgb(fg)

    cell_write(0, 0, "Brand", bg=_INK, fg=_WHITE, bold=True)
    cell_write(0, 1, "Overall", bg=_INK, fg=_WHITE, bold=True)
    for ci, d in enumerate(dims):
        cell_write(0, 2 + ci, d, bg=_INK, fg=_WHITE, bold=True, size=7)
    for ri, r in enumerate(rows, start=1):
        bg = _SELF_ROW if r.get("is_self") else _WHITE
        mark = "†" if r.get("provisional") and r["overall"]["normalized_pct"] is not None else ""
        cell_write(
            ri,
            0,
            f"{r['name']}{'  (us)' if r.get('is_self') else ''}",
            bg=bg,
            bold=bool(r.get("is_self")),
        )
        cell_write(ri, 1, f"{_fmt(r['overall']['normalized_pct'])}{mark}", bg=bg)
        for ci, d in enumerate(dims):
            sc = r.get("dimensions", {}).get(d, {}).get("score")
            if sc is None:
                cell_write(ri, 2 + ci, "", bg=_GAP_BG)  # blank, never a zero
            else:
                cell_write(ri, 2 + ci, f"{float(sc):g}", bg=bg)
    _text(
        s,
        0.7,
        6.62,
        11.9,
        0.3,
        "1–5 per dimension. Blank – not yet observed (never a zero). "
        "† provisional. Amber row – our brand.",
        size=9,
        color=_MUTED,
    )
    _footer(s, deck_title, page)


def _slide_method(
    prs: Any,
    card: dict[str, Any],
    diff: dict[str, Any] | None,
    page: int,
    deck_title: str,
) -> None:
    s = _blank(prs)
    top = _header(s, "Appendix – method", "How to read these numbers")
    dims = card.get("dimensions", [])
    th = (diff or {}).get("thresholds", {})
    if th:
        change_rule = (
            f"Material change – score move ≥ {th.get('min_score_delta', 1.0)}, "
            f"a rank change, a dimension newly scored or withdrawn, or a "
            f"coverage shift ≥ {th.get('min_coverage_delta', 0)}%. Smaller "
            "moves are suppressed."
        )
    else:
        change_rule = (
            "Material change – a full-point score move, a rank change, a "
            "dimension newly scored or withdrawn, or a large coverage shift. "
            "Smaller moves are suppressed."
        )
    items = [
        f"{len(dims)} weighted dimensions (weights sum to "
        f"{card.get('weight_total_pct', 0):g}%); each scored 1–5 from filed evidence.",
        "Overall – weighted score normalized to the weight actually measured, "
        "so a brand with partial evidence is compared on what was seen, never "
        "padded with zeros.",
        "A brand ranks only once enough of the model is scored; below that it "
        "is provisional (†) – its figure is shown, its standing withheld.",
        "Every evidence item quotes its source verbatim and was verified "
        "against the live page before it was saved; per-state evidence "
        "carries the network exit that was proven to be in that state, and "
        "storefront exhibits were captured through that same verified exit.",
        change_rule,
        "Narrative, titles, competitor observations and implications are "
        "model-written from the factual record and labelled; standings, "
        "scores and gaps are computed.",
    ]
    _bullets(
        s,
        0.7,
        top + 0.1,
        11.9,
        6.5 - top,
        items[:6],
        size=11.5,
        gap_pt=7,
        cap=260,
        accent_bullet=False,
    )
    _footer(s, deck_title, page)


def _slide_closing(prs: Any, narrative: dict[str, Any], deck_title: str) -> None:
    s = _blank(prs)
    _dark(s)
    _eyebrow(s, "Next steps", y=1.0)
    _text(s, 0.7, 1.4, 11.9, 1.0, "Where this goes next", size=30, bold=True, color=_WHITE)
    _rule(s, 2.5)
    steps = [str(b) for b in narrative.get("next_steps") or []][:4]
    if not steps:
        steps = ["Hold the cadence – re-observe on schedule and diff next cycle"]
    _bullets(s, 0.7, 2.85, 11.4, 3.4, steps, size=15, color=_DARK_BODY, gap_pt=14, cap=120)
    _judgement_note(s, str(narrative.get("source") or "facts"), dark=True)


# ── entry point ──────────────────────────────────────────────────────


def render_executive_deck(
    card: dict[str, Any],
    *,
    diff: dict[str, Any] | None,
    judged: list[dict[str, Any]],
    summary: dict[str, Any],
    gaps: list[dict[str, Any]],
    evidence_count: int,
    path: str | Path,
    screenshots: dict[str, list[dict[str, str]]] | None = None,
    offers: list[dict[str, Any]] | None = None,
    events: list[dict[str, Any]] | None = None,
    title: str = "Competitive Intelligence – Executive Briefing",
    market_label: str = "",
) -> str:
    """Write the deck. Returns the path written.

    ``summary`` is the narrative dict — ``{headline, bullets, exec, profiles,
    titles, commentary, next_steps, source}`` — model-written when a router
    exists, otherwise :func:`factual_narrative`'s computed fallback. The
    renderer treats every narrative string as untrusted copy: capped,
    en-dashed and scrubbed of internal bookkeeping. ``screenshots`` maps
    brand name → storefront exhibits ``[{path, url, observed_at}]``; slides
    that need them are skipped when they are absent, never faked.
    """
    from pptx import Presentation
    from pptx.util import Inches

    narrative = summary or {}
    shots = screenshots or {}
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    rows = card.get("rows", [])
    dims = card.get("dimensions", [])
    weights = {d["name"]: float(d.get("weight_pct") or 0.0) for d in dims}
    generated = str(card.get("generated_at") or datetime.now(UTC).isoformat())
    if diff:
        period = (
            f"Period {str(diff.get('from_generated_at', '?'))[:10]} → "
            f"{str(diff.get('to_generated_at', '?'))[:10]}"
        )
    else:
        period = f"Baseline · generated {generated[:10]}"
    basis = (
        f"{len(rows)} brands · {len(dims)} dimensions · {evidence_count:,} "
        "observed facts · every score traceable to a source URL"
    )
    deck_title = _clean(title, 70)
    src = str(narrative.get("source") or "facts")

    _slide_title(
        prs,
        title=title,
        market=market_label,
        period=period,
        basis=basis,
        generated=generated,
    )
    page = 2
    _slide_summary(prs, card, narrative, page, deck_title, events=events)
    page += 1
    _slide_glance(
        prs,
        card,
        evidence_count,
        gaps,
        (narrative.get("commentary") or {}).get("glance", ""),
        page,
        deck_title,
    )
    page += 1
    _slide_standings(prs, card, narrative, page, deck_title)
    page += 1
    _slide_versus(prs, card, narrative, page, deck_title)
    page += 1
    _slide_dimension_leaders(prs, card, narrative, page, deck_title)
    page += 1

    # The offers on the table — one row per brand, then the promotions pages
    # photographed, captioned with the offer they show.
    offer_rows = [o for o in (offers or []) if o.get("welcome") or o.get("ongoing")]
    if offer_rows:
        _slide_offers(prs, offer_rows, narrative, page, deck_title)
        page += 1
        promo_items: list[tuple[str, dict[str, str], bool]] = []
        captions: dict[str, str] = {}
        for o in offer_rows:
            promo_shot = o.get("exhibit")
            if isinstance(promo_shot, dict) and promo_shot.get("path"):
                promo_items.append((str(o["brand"]), promo_shot, bool(o.get("is_self"))))
                captions[str(o["brand"])] = str(o.get("welcome") or o.get("ongoing") or "")
        promo_items = promo_items[:8]
        pbatches = [promo_items[i : i + 2] for i in range(0, len(promo_items), 2)]
        exhibits_panel = (narrative.get("slides") or {}).get("offers") or {}
        for bi, batch in enumerate(pbatches):
            _slide_exhibits(
                prs,
                batch,
                bi + 1,
                len(pbatches),
                page,
                deck_title,
                eyebrow="Offers and promotions – exhibits",
                title="The offers as the visitor sees them",
                captions=captions,
                commentary=(exhibits_panel.get("observations") or [""])[0] if bi == 0 else "",
            )
            page += 1

    # Competitor deep dives — one slide per profiled brand, in the model's
    # order (the ranked leader first). Only brands that exist in the card.
    by_name = {r["name"]: r for r in rows}
    for profile in (narrative.get("profiles") or [])[:8]:
        row = by_name.get(str(profile.get("brand") or ""))
        if row is None:
            continue
        _slide_profile(
            prs,
            row,
            profile,
            weights,
            shots.get(row["name"], []),
            page,
            deck_title,
        )
        page += 1

    # Storefront exhibits — ranked order, our brand first when captured.
    exhibit_items: list[tuple[str, dict[str, str], bool]] = []
    ordered_rows = sorted(
        rows,
        key=lambda r: (
            not r.get("is_self"),
            r.get("rank") if r.get("rank") is not None else 99,
        ),
    )
    for r in ordered_rows:
        for shot in shots.get(r["name"], [])[:1]:
            exhibit_items.append((r["name"], shot, bool(r.get("is_self"))))
    exhibit_items = exhibit_items[:14]
    batches = [exhibit_items[i : i + 2] for i in range(0, len(exhibit_items), 2)]
    store_panel = (narrative.get("slides") or {}).get("exhibits") or {}
    store_obs = [str(o) for o in (store_panel.get("observations") or []) if str(o).strip()]
    for bi, batch in enumerate(batches):
        _slide_exhibits(
            prs,
            batch,
            bi + 1,
            len(batches),
            page,
            deck_title,
            commentary=store_obs[bi] if bi < len(store_obs) else "",
        )
        page += 1

    _slide_changes(prs, diff, narrative, page, deck_title, events=events)
    page += 1
    page += _slide_implications(
        prs,
        judged,
        int((diff or {}).get("material_count", 0)),
        src,
        page,
        deck_title,
    )
    _slide_decisions(
        prs,
        judged,
        src,
        int((diff or {}).get("material_count", 0)),
        page,
        deck_title,
    )
    page += 1
    _slide_evidence(prs, card, gaps, evidence_count, narrative, page, deck_title)
    page += 1
    _slide_heatmap(prs, card, page, deck_title)
    page += 1
    _slide_method(prs, card, diff, page, deck_title)
    page += 1
    _slide_closing(prs, narrative, deck_title)

    out = Path(path).expanduser()
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    logger.info("Executive deck written: %s (%d slides)", out, len(prs.slides))
    return str(out)
