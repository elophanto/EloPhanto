"""Executive narrative rules and search-driven source expansion.

Two complaints from the first customer-facing pack (2026-08-15):

1. "It mentions some SHA which makes zero sense. It has no design." The deck
   carried internal bookkeeping and read as a dump. The narrative prompt now
   bans machinery talk outright, and the renderer scrubs whatever slips
   through — both layers tested here.

2. "For some of the casinos it didn't spot info. It should research
   properly." The brand's own site was the only source, and for brands
   behind JS shells and bot walls that meant "not publicly observable" while
   reviews and help centers carried the facts in plain text. Expansion
   searches for third-party pages, reads them through the same verified
   exit, and holds them to the same verbatim-excerpt gate.
"""

from __future__ import annotations

from dataclasses import dataclass
from typing import Any

import pytest

import core.watch_observe as wo
from core.watch import WatchManager
from core.watch_deck import _clean
from tools.watch.tools import _DECK_NARRATIVE_SYSTEM


class TestNarrativeContract:
    def test_bans_internal_bookkeeping_by_name(self) -> None:
        text = _DECK_NARRATIVE_SYSTEM.lower()
        for term in ("hashes", "sha", "file paths", "manifests", "checkpoints",
                     "run ids"):
            assert term in text, f"prompt must ban {term!r}"
        assert "the room hears about the market, not about the machinery" in text

    def test_demands_action_titles(self) -> None:
        text = _DECK_NARRATIVE_SYSTEM.lower()
        assert "a sentence someone could disagree with" in text
        assert "never a label" in text

    def test_demands_en_dashes(self) -> None:
        assert "en dashes (–), never em dashes (—)" in _DECK_NARRATIVE_SYSTEM


class TestRendererScrubber:
    """The prompt is the first layer; the scrubber catches what slips."""

    def test_strips_hashes(self) -> None:
        assert "d933a396737d1868" not in _clean(
            "Verified via d933a396737d1868 capture"
        )

    def test_strips_bookkeeping_words(self) -> None:
        out = _clean("Per the SHA-256 freeze receipt and manifest, scores held.")
        assert "SHA" not in out and "manifest" not in out.lower()
        assert "freeze receipt" not in out.lower()

    def test_swaps_em_dashes(self) -> None:
        assert "—" not in _clean("Leads — but narrowly")
        assert "–" in _clean("Leads — but narrowly")

    def test_leaves_ordinary_numbers_alone(self) -> None:
        out = _clean("Revenue of 1000000 across 14 brands in 2026")
        assert "1000000" in out and "14" in out and "2026" in out

    def test_caps(self) -> None:
        assert len(_clean("x" * 999, 120)) == 120


class TestExpansionPrimitives:
    def test_queries_group_missing_dimensions(self) -> None:
        qs = wo.expansion_queries(
            "McLuck",
            ["Payment options and limits", "AMOE policy", "Loyalty programme",
             "KYC strategy"],
        )
        assert qs and all(q.startswith('"McLuck"') for q in qs)
        assert len(qs) <= 3
        assert "Payment" in qs[0]

    def test_no_missing_dimensions_no_queries(self) -> None:
        assert wo.expansion_queries("McLuck", []) == []

    def test_url_picker_dedupes_hosts_and_skips_already_fetched(self) -> None:
        results = [
            {"url": "https://reviews.example/mcluck", "title": "", "snippet": ""},
            {"url": "https://reviews.example/mcluck-2", "title": "", "snippet": ""},
            {"url": "https://www.mcluck.com/", "title": "", "snippet": ""},
            {"url": "ftp://bad.example/x", "title": "", "snippet": ""},
            {"url": "https://help.example/article#frag", "title": "", "snippet": ""},
        ]
        picked = wo.pick_expansion_urls(
            results, already_fetched={"https://www.mcluck.com/"}, limit=4
        )
        assert picked == [
            "https://reviews.example/mcluck",
            "https://help.example/article",
        ]


@dataclass
class _Resp:
    content: str


class _Router:
    """Extracts a claim only for the dimension the page supports."""

    async def complete(self, *, messages: list[dict[str, str]], **_kw: Any) -> _Resp:
        return _Resp(
            '{"claims": [{"dimension": "Payments", "subcriterion": "methods", '
            '"claim": "McLuck supports Visa and ACH withdrawals", '
            '"value_text": "Visa, ACH", '
            '"excerpt": "McLuck supports Visa and ACH withdrawals for players"}]}'
        )


class _Vault:
    def __init__(self, key: str | None) -> None:
        self._key = key

    def get(self, name: str) -> str | None:
        return self._key if name == "search_sh_api_key" else None


@pytest.fixture
async def wm(tmp_path):
    from core.database import Database

    db = Database(str(tmp_path / "exp.db"))
    await db.initialize()
    yield WatchManager(db)
    await db.close()


async def _analyze(wm, monkeypatch, *, vault_key, search_results):
    from tools.watch.tools import WatchAnalyzeTool

    await wm.upsert_dimension(
        name="Game portfolio", company_id="c1", weight_pct=50,
        subcriteria=[{"name": "range", "weight_pct": 100}],
    )
    await wm.upsert_dimension(
        name="Payments", company_id="c1", weight_pct=50,
        subcriteria=[{"name": "methods", "weight_pct": 100}],
    )
    await wm.add_subject(name="McLuck", company_id="c1", url="https://www.mcluck.com")

    # The brand's own site talks about games only — Payments stays silent.
    async def fake_collect(start_url, **kw):
        return [{
            "url": start_url,
            "text": "Play over 1,000 casino-style games at McLuck today " * 20,
            "error": None, "method": "http",
        }]

    site_claims = (
        '{"claims": [{"dimension": "Game portfolio", "subcriterion": "range", '
        '"claim": "McLuck offers over 1,000 games", "value_text": "1,000", '
        '"excerpt": "Play over 1,000 casino-style games at McLuck today"}]}'
    )

    class SiteThenExpandRouter:
        def __init__(self) -> None:
            self.calls = 0

        async def complete(self, **_kw: Any) -> _Resp:
            self.calls += 1
            if self.calls == 1:
                return _Resp(site_claims)
            _kw.pop("messages", None)
            return await _Router().complete(messages=[], **_kw)

    async def fake_search(query, *, api_key, **kw):
        assert api_key == "sk-test"
        return search_results

    async def fake_fetch(url, **kw):
        return ("McLuck supports Visa and ACH withdrawals for players " * 10,
                None, "http")

    monkeypatch.setattr(wo, "collect_pages", fake_collect)
    monkeypatch.setattr(wo, "search_web", fake_search)
    # The tool imports this inside execute(), so patching the source module
    # is what its function-scope import resolves against.
    monkeypatch.setattr(wo, "fetch_page_best_effort", fake_fetch)

    t = WatchAnalyzeTool()
    t._watch_manager = wm
    t._router = SiteThenExpandRouter()
    t._config = None
    t._vault = _Vault(vault_key)
    return await t.execute({"subject": "McLuck", "company_id": "c1",
                            "save": False, "deck": False})


class TestExpansionInAnalyze:
    @pytest.mark.asyncio
    async def test_silent_dimensions_get_third_party_evidence(
        self, wm, monkeypatch
    ) -> None:
        res = await _analyze(
            wm, monkeypatch, vault_key="sk-test",
            search_results=[{"url": "https://reviews.example/mcluck",
                             "title": "McLuck review", "snippet": ""}],
        )
        assert res.success, res.error
        exp = res.data["source_expansion"]
        assert exp["attempted"] is True
        assert exp["missing_dimensions"] == ["Payments"]
        assert exp["evidence_written"] == 1

        rows = await wm.list_evidence("c1")
        third_party = [r for r in rows if r.source_type == "third_party"]
        assert len(third_party) == 1
        assert third_party[0].source_url == "https://reviews.example/mcluck"
        assert third_party[0].confidence == "low"
        assert "Visa and ACH" in third_party[0].claim

    @pytest.mark.asyncio
    async def test_no_vault_key_skips_and_says_so(self, wm, monkeypatch) -> None:
        res = await _analyze(wm, monkeypatch, vault_key=None, search_results=[])
        assert res.success, res.error
        exp = res.data["source_expansion"]
        assert exp["attempted"] is False
        assert "search_sh_api_key" in exp["note"]
        rows = await wm.list_evidence("c1")
        assert all(r.source_type == "site" for r in rows)

    @pytest.mark.asyncio
    async def test_covered_dimensions_are_not_searched(self, wm, monkeypatch) -> None:
        """Expansion targets only what the site left silent."""
        res = await _analyze(
            wm, monkeypatch, vault_key="sk-test",
            search_results=[{"url": "https://reviews.example/mcluck",
                             "title": "", "snippet": ""}],
        )
        exp = res.data["source_expansion"]
        assert "Game portfolio" not in exp["missing_dimensions"]


class TestUnreadableSiteStillGetsResearched:
    """The exact Chumba failure, 2026-08-15 19:47:

        could not read any page for Chumba Casino:
          ['browser exit not verified in FL (AL) — page dropped …']

    The site was bot-walled, Chrome's exit landed in Alabama, exit
    verification (correctly) dropped the page — and the early return then
    skipped source expansion entirely. The brands that most need third-party
    research were exactly the ones that never got it.
    """

    @pytest.mark.asyncio
    async def test_expansion_rescues_a_zero_page_brand(self, wm, monkeypatch) -> None:
        from tools.watch.tools import WatchAnalyzeTool

        await wm.upsert_dimension(
            name="Payments", company_id="c1", weight_pct=100,
            subcriteria=[{"name": "methods", "weight_pct": 100}],
        )
        await wm.add_subject(name="Chumba Casino", company_id="c1",
                             url="https://www.chumbacasino.com")

        async def unreadable(start_url, **kw):
            return [{"url": start_url, "text": "",
                     "error": "browser exit not verified in FL (AL)",
                     "method": "browser"}]

        async def fake_search(query, *, api_key, **kw):
            return [{"url": "https://reviews.example/chumba", "title": "", "snippet": ""}]

        async def fake_fetch(url, **kw):
            return ("Chumba supports Visa and ACH withdrawals for players " * 10,
                    None, "http")

        monkeypatch.setattr(wo, "collect_pages", unreadable)
        monkeypatch.setattr(wo, "search_web", fake_search)
        monkeypatch.setattr(wo, "fetch_page_best_effort", fake_fetch)

        class R:
            async def complete(self, **_kw: Any) -> _Resp:
                return _Resp(
                    '{"claims": [{"dimension": "Payments", "subcriterion": "methods", '
                    '"claim": "Chumba supports Visa and ACH withdrawals", '
                    '"value_text": "Visa, ACH", '
                    '"excerpt": "Chumba supports Visa and ACH withdrawals for players"}]}'
                )

        t = WatchAnalyzeTool()
        t._watch_manager = wm
        t._router = R()
        t._config = None
        t._vault = _Vault("sk-test")
        res = await t.execute({"subject": "Chumba Casino", "company_id": "c1",
                               "save": False, "deck": False})

        assert res.success, res.error
        assert res.data["source_expansion"]["evidence_written"] == 1
        rows = await wm.list_evidence("c1")
        assert rows and rows[0].source_type == "third_party"

    @pytest.mark.asyncio
    async def test_zero_pages_and_no_expansion_is_still_an_honest_error(
        self, wm, monkeypatch
    ) -> None:
        from tools.watch.tools import WatchAnalyzeTool

        await wm.upsert_dimension(
            name="Payments", company_id="c1", weight_pct=100,
            subcriteria=[{"name": "methods", "weight_pct": 100}],
        )
        await wm.add_subject(name="Walled", company_id="c1", url="https://w.example")

        async def unreadable(start_url, **kw):
            return [{"url": start_url, "text": "", "error": "HTTP 403",
                     "method": "http"}]

        monkeypatch.setattr(wo, "collect_pages", unreadable)

        t = WatchAnalyzeTool()
        t._watch_manager = wm
        t._router = _Router()
        t._config = None
        t._vault = _Vault(None)  # no search key
        res = await t.execute({"subject": "Walled", "company_id": "c1",
                               "save": False, "deck": False})

        assert not res.success
        assert "could not read any page" in res.error
        assert "search_sh_api_key" in res.error  # says why expansion was skipped


class TestChromeSessionPinning:
    def test_pin_password_appends_session_to_geo_passwords(self) -> None:
        out = wo.pin_password("pw_country-us_state-florida", session="ab12cd34")
        assert out == "pw_country-us_state-florida_session-ab12cd34_lifetime-30m"

    def test_pin_password_is_idempotent_and_ignores_plain_passwords(self) -> None:
        pinned = wo.pin_password("pw_state-florida", session="ab12cd34")
        assert wo.pin_password(pinned) == pinned
        assert wo.pin_password("plainsecret") == "plainsecret"

    def test_agent_pins_the_browser_password(self) -> None:
        """Chrome must hold one exit per session window — an unpinned
        rotating credential exits anywhere in the country (Alabama,
        2026-08-15), and every browser-escalated page of a state-stamped
        run is then correctly dropped."""
        import inspect

        import core.agent as agent_mod

        src = inspect.getsource(agent_mod)
        assert "pin_password(" in src, (
            "browser proxy_password must be session-pinned at wiring time"
        )


class TestBrowserFirstLinkDiscovery:
    """'It's trying pages URLs instead of going through the site actually.'

    Link discovery parsed the raw HTTP response — a JS-app site serves an
    empty shell there (big HTML, zero anchors), so for exactly the walled
    brands the crawler concluded the site had one page. The rendered DOM
    knows the nav the shell hides; when raw discovery finds nothing and
    Chrome is available, discovery now uses what Chrome actually renders.
    """

    @pytest.mark.asyncio
    async def test_shell_homepage_falls_back_to_rendered_dom(self, monkeypatch) -> None:
        fetched: list[str] = []

        async def fake_raw(url, **kw):
            return "<html><div id=root></div><script src=/app.js></script></html>"

        async def fake_browser_html(bm, url):
            return ('<a href="/sweeps-rules">Sweepstakes Rules</a>'
                    '<a href="/promotions">Promotions</a>')

        async def fake_fetch(url, **kw):
            fetched.append(url)
            return ("page text " * 100, None, "browser")

        monkeypatch.setattr(wo, "_fetch_raw_html", fake_raw)
        monkeypatch.setattr(wo, "_browser_html", fake_browser_html)
        monkeypatch.setattr(wo, "fetch_page_best_effort", fake_fetch)

        pages = await wo.collect_pages(
            "https://walled.example", browser_manager=object(), max_pages=4
        )
        assert [p["url"] for p in pages][1:] == [
            "https://walled.example/sweeps-rules",
            "https://walled.example/promotions",
        ]

    @pytest.mark.asyncio
    async def test_server_rendered_sites_never_touch_the_browser(
        self, monkeypatch
    ) -> None:
        """Raw discovery that works must stay the cheap path."""
        called = {"browser": 0}

        async def fake_raw(url, **kw):
            return '<a href="/terms">Terms of use</a>'

        async def fake_browser_html(bm, url):
            called["browser"] += 1
            return ""

        async def fake_fetch(url, **kw):
            return ("page text " * 100, None, "http")

        monkeypatch.setattr(wo, "_fetch_raw_html", fake_raw)
        monkeypatch.setattr(wo, "_browser_html", fake_browser_html)
        monkeypatch.setattr(wo, "fetch_page_best_effort", fake_fetch)

        pages = await wo.collect_pages(
            "https://plain.example", browser_manager=object(), max_pages=3
        )
        assert called["browser"] == 0
        assert len(pages) == 2  # homepage + terms


class _CapturingBrowser:
    """Fake browser manager: navigate succeeds, browser_capture writes a
    real file at the requested path — the bridge contract, minus Chrome."""

    def __init__(self, *, fail: bool = False) -> None:
        self.fail = fail
        self.calls: list[tuple[str, dict[str, Any]]] = []

    async def call_tool(self, name: str, params: dict[str, Any]) -> dict[str, Any]:
        self.calls.append((name, params))
        if name == "browser_navigate":
            return {"success": True, "url": params.get("url")}
        if name == "browser_capture":
            if self.fail:
                return {"success": False, "error": "no page"}
            from pathlib import Path

            out = Path(str(params["path"]))
            out.parent.mkdir(parents=True, exist_ok=True)
            out.write_bytes(b"\xff\xd8\xff\xe0 fake jpeg")
            return {"success": True, "path": str(out), "bytes": 14}
        return {"success": True}


class TestCapturePageScreenshot:
    @pytest.mark.asyncio
    async def test_navigates_and_writes_the_file(self, tmp_path) -> None:
        from core.watch_observe import capture_page_screenshot

        bm = _CapturingBrowser()
        out = tmp_path / "shots" / "home.jpg"
        got = await capture_page_screenshot(bm, "https://x.example", str(out))
        assert got == str(out) and out.exists()
        # navigate → consent dismissal (the browser playbook's click_text
        # recipe, best-effort) → capture. The capture is always last.
        assert bm.calls[0][0] == "browser_navigate"
        assert bm.calls[-1][0] == "browser_capture"
        assert any(c[0] == "browser_click_text" for c in bm.calls)

    @pytest.mark.asyncio
    async def test_failure_returns_empty_never_raises(self, tmp_path) -> None:
        from core.watch_observe import capture_page_screenshot

        got = await capture_page_screenshot(
            _CapturingBrowser(fail=True), "https://x.example",
            str(tmp_path / "no.jpg"),
        )
        assert got == ""
        assert await capture_page_screenshot(None, "https://x", "/tmp/x.jpg") == ""

    @pytest.mark.asyncio
    async def test_a_font_stall_gets_one_more_shot(self, tmp_path) -> None:
        """High 5, 2026-08-16: page.screenshot timed out 'waiting for fonts
        to load' on the only exhibit page. The second shot lands."""
        from core.watch_observe import capture_page_screenshot

        class _Stalls(_CapturingBrowser):
            def __init__(self) -> None:
                super().__init__()
                self.shots = 0

            async def call_tool(self, name, params):
                if name == "browser_capture":
                    self.shots += 1
                    if self.shots == 1:
                        self.calls.append((name, params))
                        return {"success": False,
                                "error": "page.screenshot: Timeout 30000ms exceeded. waiting for fonts to load..."}
                return await super().call_tool(name, params)

        bm = _Stalls()
        out = tmp_path / "h5" / "home.jpg"
        assert await capture_page_screenshot(bm, "https://h5.example", str(out)) == str(out)
        assert bm.shots == 2 and out.exists()
        # consent is dismissed again between the stalled shot and the retry:
        # a bar that slid in during the 30s stall must not be in the picture
        names = [c[0] for c in bm.calls]
        first_shot = names.index("browser_capture")
        assert "browser_click_text" in names[first_shot + 1 : len(names) - 1]
        # and a browser that fails twice is a gap, not an exception
        assert await capture_page_screenshot(_CapturingBrowser(fail=True), "https://x", str(tmp_path / "n.jpg")) == ""

    def test_filenames_are_dated_slugs(self) -> None:
        from core.watch_observe import screenshot_filename

        name = screenshot_filename(
            "https://www.mcluck.com/Promotions/Daily-Wheel?x=1", when="20260815"
        )
        assert name == "20260815-promotions-daily-wheel.jpg"
        assert screenshot_filename("https://x.example/", when="20260815") == (
            "20260815-home.jpg"
        )


class TestStorefrontExhibitsInAnalyze:
    """watch_analyze files clean storefront captures and stamps the evidence
    rows from that page with the exhibit path."""

    async def _run(self, wm, monkeypatch, tmp_path, *, browser) -> Any:
        from types import SimpleNamespace

        from tools.watch.tools import WatchAnalyzeTool

        await wm.upsert_dimension(
            name="Game portfolio", company_id="c1", weight_pct=100,
            subcriteria=[{"name": "range", "weight_pct": 100}],
        )
        await wm.add_subject(
            name="McLuck", company_id="c1", url="https://www.mcluck.com"
        )

        async def fake_collect(start_url, **kw):
            return [{
                "url": start_url,
                "text": "Play over 1,000 casino-style games at McLuck today " * 20,
                "error": None, "method": "http",
            }]

        monkeypatch.setattr(wo, "collect_pages", fake_collect)

        class R:
            async def complete(self, **_kw: Any) -> _Resp:
                return _Resp(
                    '{"claims": [{"dimension": "Game portfolio", '
                    '"subcriterion": "range", '
                    '"claim": "McLuck offers over 1,000 games", '
                    '"value_text": "1,000", '
                    '"excerpt": "Play over 1,000 casino-style games at McLuck today"}]}'
                )

        t = WatchAnalyzeTool()
        t._watch_manager = wm
        t._router = R()
        t._vault = _Vault(None)
        t._browser_manager = browser
        t._config = SimpleNamespace(
            workspace=str(tmp_path / "ws"), project_root=tmp_path, proxy=None
        )
        return await t.execute({
            "subject": "McLuck", "company_id": "c1",
            "save": False, "deck": False, "expand_sources": False,
        })

    @pytest.mark.asyncio
    async def test_captures_and_stamps_evidence_rows(
        self, wm, monkeypatch, tmp_path
    ) -> None:
        bm = _CapturingBrowser()
        res = await self._run(wm, monkeypatch, tmp_path, browser=bm)
        assert res.success, res.error
        shots = res.data["screenshots"]
        assert shots["captured"] == 1
        path = shots["paths"][0]
        assert "watch-screenshots/mcluck/" in path and path.endswith(".jpg")
        from pathlib import Path

        assert Path(path).exists()
        rows = await wm.list_evidence("c1")
        assert rows and all(r.screenshot_path == path for r in rows)

    @pytest.mark.asyncio
    async def test_no_browser_means_no_exhibits_and_no_error(
        self, wm, monkeypatch, tmp_path
    ) -> None:
        res = await self._run(wm, monkeypatch, tmp_path, browser=None)
        assert res.success, res.error
        # A missing browser is a VISIBLE skip now, never a silent None —
        # the silent skip is how a whole run shipped without exhibits.
        assert res.data["screenshots"]["captured"] == 0
        assert "no browser" in res.data["screenshots"]["note"]
        rows = await wm.list_evidence("c1")
        assert rows and all(r.screenshot_path == "" for r in rows)

    @pytest.mark.asyncio
    async def test_capture_failure_is_noted_not_fatal(
        self, wm, monkeypatch, tmp_path
    ) -> None:
        res = await self._run(
            wm, monkeypatch, tmp_path, browser=_CapturingBrowser(fail=True)
        )
        assert res.success, res.error
        assert res.data["screenshots"]["captured"] == 0
        assert "capture failed" in res.data["screenshots"]["note"]

    @pytest.mark.asyncio
    async def test_geo_run_refuses_exhibits_on_unverified_browser_exit(
        self, wm, monkeypatch, tmp_path
    ) -> None:
        """A Florida-stamped run must not exhibit an Alabama storefront."""
        from types import SimpleNamespace

        from tools.watch.tools import WatchAnalyzeTool

        await wm.upsert_dimension(
            name="Game portfolio", company_id="c1", weight_pct=100,
            subcriteria=[{"name": "range", "weight_pct": 100}],
        )
        await wm.add_subject(
            name="McLuck", company_id="c1", url="https://www.mcluck.com"
        )

        async def fake_collect(start_url, **kw):
            return [{
                "url": start_url,
                "text": "Play over 1,000 casino-style games at McLuck today " * 20,
                "error": None, "method": "http",
            }]

        async def fake_verify_exit(proxy_url, state, **kw):
            return True, proxy_url, {"ip": "1.2.3.4", "state_code": "FL"}

        async def fake_verify_browser(bm, state, **kw):
            return False, {"state_code": "AL", "ip": "5.6.7.8"}

        monkeypatch.setattr(wo, "collect_pages", fake_collect)
        monkeypatch.setattr(wo, "verify_exit_state", fake_verify_exit)
        monkeypatch.setattr(wo, "verify_browser_exit", fake_verify_browser)

        class R:
            async def complete(self, **_kw: Any) -> _Resp:
                return _Resp(
                    '{"claims": [{"dimension": "Game portfolio", '
                    '"subcriterion": "range", '
                    '"claim": "McLuck offers over 1,000 games", '
                    '"value_text": "1,000", '
                    '"excerpt": "Play over 1,000 casino-style games at McLuck today"}]}'
                )

        bm = _CapturingBrowser()
        t = WatchAnalyzeTool()
        t._watch_manager = wm
        t._router = R()
        t._vault = _Vault(None)
        t._browser_manager = bm
        t._config = SimpleNamespace(
            workspace=str(tmp_path / "ws"), project_root=tmp_path,
            proxy=SimpleNamespace(
                request_proxy_url=lambda state: "http://sticky.example:1"
            ),
        )
        res = await t.execute({
            "subject": "McLuck", "company_id": "c1", "geo_state": "FL",
            "save": False, "deck": False, "expand_sources": False,
        })
        assert res.success, res.error
        shots = res.data["screenshots"]
        assert shots["captured"] == 0
        assert "not verified in FL" in shots["note"]
        assert not any(c[0] == "browser_capture" for c in bm.calls)
        rows = await wm.list_evidence("c1")
        assert all(r.screenshot_path == "" for r in rows)


class TestCollectExhibits:
    """brand → exhibits mapping for the deck: register first, workspace
    fallback for brands whose shots never landed on a row."""

    def test_register_paths_win_and_missing_files_are_dropped(self, tmp_path) -> None:
        from tools.watch.tools import _collect_exhibits

        real = tmp_path / "real.jpg"
        real.write_bytes(b"x")
        evidence = [
            {"subject": "A", "screenshot_path": str(real),
             "source_url": "https://a.example", "observed_at": "2026-08-15"},
            {"subject": "A", "screenshot_path": str(real),  # duplicate
             "source_url": "https://a.example/2", "observed_at": "2026-08-15"},
            {"subject": "B", "screenshot_path": str(tmp_path / "gone.jpg"),
             "source_url": "https://b.example", "observed_at": "2026-08-15"},
            {"subject": "C", "screenshot_path": "",
             "source_url": "https://c.example", "observed_at": "2026-08-15"},
        ]
        out = _collect_exhibits(evidence, [{"name": "A"}, {"name": "B"}], None)
        assert list(out) == ["A"]
        assert out["A"] == [{
            "path": str(real), "url": "https://a.example",
            "observed_at": "2026-08-15", "kind": "home",
        }]

    def test_workspace_fallback_covers_walled_brands(self, tmp_path) -> None:
        from types import SimpleNamespace

        from tools.watch.tools import _collect_exhibits

        d = tmp_path / "ws" / "watch-screenshots" / "walled-brand"
        d.mkdir(parents=True)
        (d / "20260814-home.jpg").write_bytes(b"x")
        (d / "20260815-home.jpg").write_bytes(b"x")
        cfg = SimpleNamespace(workspace=str(tmp_path / "ws"), project_root=tmp_path)
        out = _collect_exhibits([], [{"name": "Walled Brand"}], cfg)
        from pathlib import Path

        assert [Path(p["path"]).name for p in out["Walled Brand"]] == [
            "20260815-home.jpg", "20260814-home.jpg",
        ]


class TestVaultReachesTheWatchOrgan:
    """Regression: the vault injection list omitted watch_analyze, so
    expansion always reported 'no search_sh_api_key in vault' while the key
    sat in the vault. The injection list is source of truth — pin it."""

    def test_watch_analyze_is_in_the_vault_injection_list(self) -> None:
        import inspect

        from core.agent import Agent

        src = inspect.getsource(Agent._inject_vault_deps)
        assert '"watch_analyze"' in src, (
            "watch_analyze must receive the vault, or source expansion "
            "silently loses its search key"
        )

    def test_watch_analyze_declares_the_vault_slot(self) -> None:
        from tools.watch.tools import WatchAnalyzeTool

        t = WatchAnalyzeTool()
        assert hasattr(t, "_vault") and t._vault is None


class TestExhibitsAreNotAModelChoice:
    """Regression 2026-08-15: `screenshots` spent one evening in the input
    schema. The first goal-driven run under timeout pressure passed
    screenshots=false to save time and shipped a pack with no exhibits.
    Exhibits are collection, not an option the model weighs against the
    clock — the kwarg stays for programmatic callers, the schema stays
    silent."""

    def test_schema_does_not_expose_screenshots(self) -> None:
        from tools.watch.tools import WatchAnalyzeTool

        assert "screenshots" not in WatchAnalyzeTool().input_schema["properties"]

    @pytest.mark.asyncio
    async def test_kwarg_still_works_for_programmatic_callers(
        self, wm, monkeypatch, tmp_path
    ) -> None:
        from types import SimpleNamespace

        from tools.watch.tools import WatchAnalyzeTool

        await wm.upsert_dimension(
            name="Game portfolio", company_id="c1", weight_pct=100,
            subcriteria=[{"name": "range", "weight_pct": 100}],
        )
        await wm.add_subject(
            name="McLuck", company_id="c1", url="https://www.mcluck.com"
        )

        async def fake_collect(start_url, **kw):
            return [{
                "url": start_url,
                "text": "Play over 1,000 casino-style games at McLuck today " * 20,
                "error": None, "method": "http",
            }]

        monkeypatch.setattr(wo, "collect_pages", fake_collect)

        class R:
            async def complete(self, **_kw: Any) -> _Resp:
                return _Resp(
                    '{"claims": [{"dimension": "Game portfolio", '
                    '"subcriterion": "range", '
                    '"claim": "McLuck offers over 1,000 games", '
                    '"value_text": "1,000", '
                    '"excerpt": "Play over 1,000 casino-style games at McLuck today"}]}'
                )

        bm = _CapturingBrowser()
        t = WatchAnalyzeTool()
        t._watch_manager = wm
        t._router = R()
        t._vault = _Vault(None)
        t._browser_manager = bm
        t._config = SimpleNamespace(
            workspace=str(tmp_path / "ws"), project_root=tmp_path, proxy=None
        )
        res = await t.execute({
            "subject": "McLuck", "company_id": "c1", "save": False,
            "deck": False, "expand_sources": False, "screenshots": False,
        })
        assert res.success, res.error
        assert res.data["screenshots"]["captured"] == 0
        assert res.data["screenshots"]["note"] == "disabled by caller"
        assert not any(c[0] == "browser_capture" for c in bm.calls)


class TestAutonomousRunsCannotGrowTheRegister:
    """Regression 2026-08-15: a goal executor 'completing the canon' added
    Fortune Coins and Global Poker to the customer's 14-brand register.
    Auto-add on url= is an operator convenience; in goal/mind/scheduled
    contexts the register is read-only canon and the tool refuses."""

    @pytest.mark.asyncio
    async def test_goal_context_cannot_add_a_brand(self, wm) -> None:
        from types import SimpleNamespace

        from core.execution_context import TaskSource, execution_context
        from tools.watch.tools import WatchAnalyzeTool

        t = WatchAnalyzeTool()
        t._watch_manager = wm
        t._router = SimpleNamespace()  # never reached
        with execution_context(source=TaskSource.GOAL):
            res = await t.execute({
                "subject": "Invented Casino", "company_id": "c1",
                "url": "https://invented.example",
            })
        assert not res.success
        assert "register is canon" in res.error
        assert await wm.get_subject_by_name("Invented Casino", "c1") is None

    @pytest.mark.asyncio
    async def test_operator_context_still_adds(self, wm, monkeypatch) -> None:

        from tools.watch.tools import WatchAnalyzeTool

        async def fake_collect(start_url, **kw):
            return [{"url": start_url, "text": "words " * 100,
                     "error": None, "method": "http"}]

        monkeypatch.setattr(wo, "collect_pages", fake_collect)
        await wm.upsert_dimension(
            name="Promos", company_id="c1", weight_pct=100,
            subcriteria=[{"name": "a", "weight_pct": 100}],
        )

        class R:
            async def complete(self, **_kw: Any) -> _Resp:
                return _Resp('{"claims": []}')

        t = WatchAnalyzeTool()
        t._watch_manager = wm
        t._router = R()
        t._vault = _Vault(None)
        t._config = None
        # Default context is USER — the operator convenience keeps working.
        await t.execute({
            "subject": "New Brand", "company_id": "c1",
            "url": "https://new.example", "save": False, "deck": False,
            "expand_sources": False,
        })
        assert await wm.get_subject_by_name("New Brand", "c1") is not None


class TestTheReceiptsFourFindings:
    """The 2026-08-15 from-scratch run's own validation receipt named four
    pipeline defects, correctly. Each is pinned here."""

    @pytest.mark.asyncio
    async def test_export_is_active_register_only_and_uncapped(self, wm) -> None:
        d = await wm.upsert_dimension(
            name="Promos", company_id="c1", weight_pct=100,
            subcriteria=[{"name": "a", "weight_pct": 100}],
        )
        keep = await wm.add_subject(name="Keep", company_id="c1", url="https://k.example")
        gone = await wm.add_subject(name="Gone", company_id="c1", url="https://g.example")
        for i in range(2_100):
            await wm.add_evidence(
                company_id="c1", subject_id=keep.subject_id,
                dimension_id=d.dimension_id, claim=f"claim {i}",
                source_url="https://k.example/p",
            )
        await wm.add_evidence(
            company_id="c1", subject_id=gone.subject_id,
            dimension_id=d.dimension_id, claim="archived brand claim",
            source_url="https://g.example/p",
        )
        await wm.archive_subject(gone.subject_id)
        rows = await wm.evidence_with_names("c1")
        # Uncapped: all 2,100 kept rows ship (the old cap was exactly 2,000).
        assert len(rows) == 2_100
        # Active-only: the archived brand's row does not.
        assert not any(r["subject"] == "Gone" for r in rows)
        # …unless asked for.
        rows_all = await wm.evidence_with_names("c1", include_archived=True)
        assert len(rows_all) == 2_101

    def test_uneven_coverage_withholds_ranks_for_the_whole_field(self) -> None:
        from core.watch import WatchDimension, WatchScore, WatchSubject, build_scorecard

        dims = [
            WatchDimension(dimension_id=f"d{i}", company_id="c1", name=f"D{i}",
                           weight_pct=10.0)
            for i in range(10)
        ]
        a = WatchSubject(subject_id="a", company_id="c1", name="Deep")
        b = WatchSubject(subject_id="b", company_id="c1", name="Shallow")

        def scores(sid, n, val, cov=100.0):
            return {
                f"d{i}": WatchScore(score_id=f"{sid}-{i}", company_id="c1",
                                    subject_id=sid, dimension_id=f"d{i}",
                                    score=val, coverage_pct=cov)
                for i in range(n)
            }

        # Both hold scores on 100% of the model weight — but Deep's evidence
        # covers 90% of what those dimensions ask and Shallow's covers 30%.
        # Being SCORED on a dimension is not being MEASURED to depth on it
        # (LuckyLand 2026-08-16: scores on 70% of weight, 27% coverage,
        # ranked #14 in a 'non-comparable' field). Judge on coverage.
        card = build_scorecard(
            [a, b], dims,
            {"a": scores("a", 10, 4, cov=90.0), "b": scores("b", 10, 5, cov=30.0)},
        )
        assert card["comparability_note"]
        assert "ranking withheld" in card["comparability_note"]
        assert all(r["rank"] is None for r in card["rows"])
        assert all(r["provisional"] for r in card["rows"])
        assert "field not comparable" in card["rows"][0]["provisional_reason"]
        # Scores are NOT suppressed — the figures still show.
        assert all(r["overall"]["normalized_pct"] is not None for r in card["rows"])

        # Bring Shallow's coverage to 70% (spread 20) and the field ranks.
        card2 = build_scorecard(
            [a, b], dims,
            {"a": scores("a", 10, 4, cov=90.0), "b": scores("b", 10, 5, cov=70.0)},
        )
        assert card2["comparability_note"] == ""
        assert [r["rank"] for r in card2["rows"]] == [1, 2]

    @pytest.mark.asyncio
    async def test_baseline_pack_makes_no_comparison(self, wm, tmp_path) -> None:
        from tools.watch.tools import WatchBoardReportTool

        d = await wm.upsert_dimension(
            name="Promos", company_id="c1", weight_pct=100,
            subcriteria=[{"name": "a", "weight_pct": 100}],
        )
        s = await wm.add_subject(name="Brand", company_id="c1", url="https://b.example")
        await wm.add_evidence(company_id="c1", subject_id=s.subject_id,
                              dimension_id=d.dimension_id, claim="x",
                              source_url="https://b.example/p")
        await wm.set_score(company_id="c1", subject_id=s.subject_id,
                           dimension_id=d.dimension_id, score=3)
        # An earlier same-day snapshot exists — the trap.
        await wm.take_snapshot("c1", label="earlier today")
        await wm.set_score(company_id="c1", subject_id=s.subject_id,
                           dimension_id=d.dimension_id, score=5)

        t = WatchBoardReportTool()
        t._watch_manager = wm
        t._router = None
        res = await t.execute({
            "company_id": "c1", "path": str(tmp_path / "r.md"),
            "baseline": True, "deck": True, "take_snapshot": False,
        })
        assert res.success, res.error
        assert res.data["material_count"] == 0
        md = (tmp_path / "r.md").read_text()
        assert "Material changes" not in md or "baseline" in md.lower()
        from pptx import Presentation

        text = "\n".join(
            sh.text_frame.text
            for sl in Presentation(res.data["deck_path"]).slides
            for sh in sl.shapes if sh.has_text_frame
        )
        assert "Baseline established" in text
        assert "material move" not in text.lower()

        # Without the flag the same state reports the move.
        res2 = await t.execute({
            "company_id": "c1", "path": str(tmp_path / "r2.md"),
            "take_snapshot": False, "deck": False,
        })
        assert res2.data["material_count"] >= 1


class TestExhibitPageRanking:
    """A storefront pack shows the storefront. 2026-08-16: four of fourteen
    exhibits were a privacy policy, a terms page, responsible-play or a deep
    slot URL because the picker took 'the first readable pages'."""

    def test_home_then_promo_then_other_never_legal(self) -> None:
        from core.watch_observe import exhibit_kind, rank_exhibit_pages

        assert exhibit_kind("https://x.example/") == "home"
        assert exhibit_kind("https://x.example/pages/privacy-policy") == "legal"
        assert exhibit_kind("https://x.example/terms-of-service") == "legal"
        assert exhibit_kind("https://x.example/responsible-play") == "legal"
        assert exhibit_kind("https://x.example/promotions") == "promo"
        assert exhibit_kind("https://x.example/games/slots/piggy-power") == "other"
        pages = [
            {"url": "https://x.example/pages/privacy-policy"},
            {"url": "https://x.example/games/slots/piggy-power"},
            {"url": "https://x.example/promotions"},
            {"url": "https://x.example/terms-of-service"},
            {"url": "https://x.example/promotions/refer-a-friend"},
        ]
        got = rank_exhibit_pages(pages, home_url="https://x.example", limit=3)
        assert got == [
            "https://x.example",
            "https://x.example/promotions",
            "https://x.example/games/slots/piggy-power",
        ]


class TestConsentDismissal:
    """The browser playbook's own recipe (click_text 'Accept All' / 'Accept'),
    reading the tool's answer properly: browser_click_text reports success
    even when it fell through to 'Home' (High 5, 2026-08-16)."""

    @pytest.mark.asyncio
    async def test_only_a_click_that_landed_on_the_label_counts(self) -> None:
        from core.watch_observe import dismiss_consent

        class _B:
            def __init__(self) -> None:
                self.calls: list[tuple[str, dict]] = []
                self.banner_up = True

            async def call_tool(self, name: str, args: dict) -> dict:
                self.calls.append((name, args))
                if name == "browser_click_text":
                    if args["text"] == "Accept" and self.banner_up:
                        self.banner_up = False
                        return {"success": True, "matchedText": "Accept", "matchedTag": "button"}
                    # the tool's fall-through: "success" on an unrelated element
                    return {"success": True, "matchedText": "Home", "matchedTag": "a"}
                if name == "browser_get_elements":
                    return {"elements": "[0]:<a>Home</a> [1]:<button>Accept</button>"}
                return {"success": True}

        b = _B()
        n = await dismiss_consent(b)
        assert n == 1
        assert not b.banner_up
        # "Accept All"/"Allow all" fall-throughs were NOT counted as hits
        clicked = [a["text"] for nm, a in b.calls if nm == "browser_click_text"]
        assert "Accept All" in clicked and "Accept" in clicked


    @pytest.mark.asyncio
    async def test_a_late_modal_gets_one_more_look(self) -> None:
        """Pulsz Bingo, 2026-08-16: nothing to click at t=0, modal on screen
        by the capture. A miss on the first look waits and looks again once."""
        from core.watch_observe import dismiss_consent

        class _B:
            def __init__(self) -> None:
                self.calls: list[tuple[str, dict]] = []
                self.waited = False
                self.banner_up = True

            async def call_tool(self, name: str, args: dict) -> dict:
                self.calls.append((name, args))
                if name == "browser_wait":
                    self.waited = True  # the modal renders during the settle
                    return {"success": True}
                if name == "browser_click_text":
                    if self.waited and self.banner_up and args["text"] == "Accept All":
                        self.banner_up = False
                        return {"success": True, "matchedText": "Accept All", "matchedTag": "button"}
                    return {"success": False, "message": "no element"}
                if name == "browser_get_elements":
                    return {"elements": "[0]:<a>Home</a>"}
                return {"success": True}

        b = _B()
        assert await dismiss_consent(b, settle_ms=(10, 20)) == 1
        assert not b.banner_up
        # a bounded number of settle waits, not an endless retry
        clean = _B()
        clean.banner_up = False
        assert await dismiss_consent(clean, settle_ms=(10, 20)) == 0
        assert sum(1 for nm, _ in clean.calls if nm == "browser_wait") == 2


class TestMarketEvents:
    """A competitor closing must reach the executive summary and the
    market-moves slide even in a baseline pack (LuckyLand, 2026-08-16)."""

    def test_closure_claim_is_an_event_and_ordinary_claims_are_not(self) -> None:
        from core.watch_deck import market_events

        ev = [
            {"subject": "LuckyLand Slots", "claim": "LuckyLand Slots is closing on September 14, 2026.",
             "value_text": "2026-09-14", "observed_at": "2026-08-16T10:04:48+00:00",
             "source_url": "https://www.luckylandslots.com"},
            {"subject": "LuckyLand Slots", "claim": "LuckyLand Slots is closing on September 14, 2026, with its games remaining available on LuckyLand Casino.",
             "observed_at": "2026-08-16T10:03:54+00:00"},
            {"subject": "Pulsz", "claim": "New players can claim a Welcome Bundle with 200% or more extra Gold Coins."},
            {"subject": "Modo", "claim": "Modo Casino is now available in Florida after launching in the state.",
             "observed_at": "2026-08-16"},
        ]
        got = market_events(ev)
        # two phrasings of the LuckyLand closure are ONE event
        assert [e["brand"] for e in got] == ["LuckyLand Slots", "Modo"]
        assert got[0]["observed_at"] == "2026-08-16" and "closing" in got[0]["claim"]

    def test_events_render_on_summary_and_market_moves(self, tmp_path) -> None:
        from pptx import Presentation

        from core.watch_deck import factual_narrative, render_executive_deck

        card = {
            "rows": [
                {"name": "Us", "is_self": True, "rank": 1,
                 "overall": {"normalized_pct": 60.0, "coverage_pct": 70.0}, "dimensions": {}},
                {"name": "LuckyLand Slots", "is_self": False, "rank": 2,
                 "overall": {"normalized_pct": 40.0, "coverage_pct": 65.0}, "dimensions": {}},
            ],
            "dimensions": [],
        }
        events = [{"brand": "LuckyLand Slots", "claim": "LuckyLand Slots is closing on September 14, 2026.",
                   "when": "", "observed_at": "2026-08-16", "url": ""}]
        out = tmp_path / "d.pptx"
        render_executive_deck(
            card, diff=None, judged=[], summary=factual_narrative(card, None, [], []),
            gaps=[], evidence_count=3, path=out, events=events,
        )
        prs = Presentation(str(out))
        texts = ["\n".join(sh.text_frame.text for sh in sl.shapes if sh.has_text_frame) for sl in prs.slides]
        assert "MARKET EVENT" in texts[1] and "closing on September 14, 2026" in texts[1]
        moves = next(t for t in texts if "Market moves" in t or "MARKET MOVES" in t.upper())
        assert "1 market event on record" in moves and "LuckyLand Slots" in moves


class TestOfferFacts:
    def test_headline_welcome_is_the_most_specific_claim_not_the_newest(self) -> None:
        from tools.watch.tools import _offer_facts

        card = {"rows": [{"name": "Pulsz Bingo", "is_self": True, "rank": None,
                          "overall": {"normalized_pct": 60.0}}]}
        ev = [  # newest first — boilerplate is newest
            {"subject": "Pulsz Bingo", "dimension": "Promotional proposition and generosity",
             "claim": "No purchase is necessary to play.", "source_url": "u1", "observed_at": "t1"},
            {"subject": "Pulsz Bingo", "dimension": "Promotional proposition and generosity",
             "claim": "New users can sign up for 5,000 free Gold Coins.", "source_url": "u2", "observed_at": "t2"},
            {"subject": "Pulsz Bingo", "dimension": "Promotional proposition and generosity",
             "claim": "The site offers daily free rewards.", "source_url": "u3", "observed_at": "t3"},
        ]
        row = _offer_facts(card, ev, {})[0]
        assert row["welcome"] == "New users can sign up for 5,000 free Gold Coins."
        assert row["ongoing"] == "The site offers daily free rewards."
        assert row["url"] == "u2"

    def test_welcome_and_ongoing_from_promo_evidence(self) -> None:
        from tools.watch.tools import _offer_facts

        card = {
            "rows": [
                {"name": "Us", "is_self": True, "rank": 2, "overall": {"normalized_pct": 60.0}},
                {"name": "Rival", "is_self": False, "rank": 1, "overall": {"normalized_pct": 70.0}},
                {"name": "Quiet", "is_self": False, "rank": 3, "overall": {"normalized_pct": 50.0}},
            ]
        }
        ev = [
            {"subject": "Rival", "dimension": "Promotional proposition and generosity",
             "claim": "New players get 200% extra Gold Coins on first purchase.",
             "value_text": "200% extra", "source_url": "https://r.example/promotions",
             "observed_at": "2026-08-16"},
            {"subject": "Rival", "dimension": "Promotional proposition and generosity",
             "claim": "Daily login bonus of free coins.", "value_text": "Daily",
             "source_url": "https://r.example", "observed_at": "2026-08-16"},
            {"subject": "Us", "dimension": "Loyalty programme and proposition",
             "claim": "Weekly wheel spin for members.", "value_text": "",
             "source_url": "https://u.example", "observed_at": "2026-08-16"},
            {"subject": "Rival", "dimension": "Game portfolio and category range",
             "claim": "1,200 games.", "value_text": "", "source_url": "", "observed_at": ""},
        ]
        exhibits = {"Rival": [{"path": "/tmp/r-promo.jpg", "url": "https://r.example/promotions", "kind": "promo"}]}
        offers = _offer_facts(card, ev, exhibits)
        assert [o["brand"] for o in offers] == ["Us", "Rival", "Quiet"]  # ours first, then rank
        rival = offers[1]
        assert rival["welcome"].startswith("New players get 200%")
        assert rival["ongoing"].startswith("Daily login bonus")
        assert rival["exhibit"]["kind"] == "promo"
        us = offers[0]
        assert us["welcome"] == "Weekly wheel spin for members."  # newest promo claim when no welcome
        assert offers[2]["welcome"] == "" and offers[2]["ongoing"] == ""  # census, honest blank
