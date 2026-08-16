"""The executive deck: the board pack as slides, with the organ's honesty
rules intact on every one of them.

Three rules are easiest to lose on a slide, so they are pinned here against
the actual .pptx, not the input dict:

* unscored is blank, never zero — in the heatmap and the standings chart;
* provisional brands are listed, not ranked — a bar is trusted without its
  footnote, so a brand scored on 15% of the model must not lead the chart;
* fact and judgement are labelled — the summary slide says whether a model
  wrote it, and "no decision required" never masquerades for "not evaluated".
"""

from __future__ import annotations

from dataclasses import dataclass
from typing import Any

import pytest

from core.watch import WatchManager


@dataclass
class _Resp:
    content: str


class _DeckRouter:
    """Answers the judgement and the summary prompts with valid JSON."""

    def __init__(self) -> None:
        self.calls: list[str] = []

    async def complete(self, *, messages: list[dict[str, str]], **_kw: Any) -> _Resp:
        system = messages[0]["content"]
        is_narrative = "words for an executive competitor" in system
        self.calls.append("summary" if is_narrative else "judge")
        if is_narrative:
            return _Resp(
                '{"headline": "We lead, but Rival One is closing on promos.", '
                '"bullets": ["OurBrand ranks #1 at 75.0.", '
                '"Rival One dropped a full point on Promos.", '
                '"Two dimensions remain unobserved for Rival Two."], '
                '"exec": {"findings": ["Rival One cut its welcome offer this period."], '
                '"threats": ["Rival One still out-promos us on Promos."], '
                '"watch": ["Watch Rival Two – Trust newly published."], '
                '"by_dimension": [{"dimension": "Promos", "observations": '
                '["Rival One leads promos with a 200% welcome bundle.", '
                '"OurBrand trails by two points on Promos."]}, '
                '{"dimension": "Not A Dimension", "observations": ["dropped"]}], '
                '"recommendation": "Match the welcome bundle before Q4 or cede the promos lead.", '
                '"actions": ["Approve a Q4 welcome bundle", "Refresh Rival Two next cycle"]}, '
                '"titles": {"standings": "OurBrand leads a thin field"}, '
                '"commentary": {"standings": "Amber is us – the lead is real but narrow."}, '
                '"slides": {"standings": {"observations": '
                '["OurBrand leads; Rival One is one point back."], '
                '"implications": ["The lead is narrow – one promo cycle could flip it."]}, '
                '"offers": {"observations": '
                '["Rival One is the only brand stating a percentage welcome offer."], '
                '"implications": ["Our welcome offer is unstated on the pages read."]}}, '
                '"profiles": [{"brand": "Rival One", "title": "closing fast on promotions", '
                '"observations": ["Runs a two-tier welcome offer with daily login bonuses."], '
                '"implications": ["Our promo calendar needs a counter before Q4."]}, '
                '{"brand": "Rival Two", "title": "newly transparent on licensing", '
                '"observations": ["Publishes licensing and responsible-play pages."], '
                '"implications": ["Table stakes are rising on trust."]}, '
                '{"brand": "Not A Brand", "title": "hallucinated", '
                '"observations": ["should be dropped"], "implications": []}], '
                '"next_steps": ["Collect the unobserved brands"]}'
            )
        return _Resp(
            '{"items": [{"subject": "Rival One", "change": "Promos 5 -> 2", '
            '"implication": "Their welcome offer weakened.", '
            '"recommendation": "Hold our promo spend.", '
            '"classification": "no_regret", "decision_required": "none"}, '
            '{"subject": "Rival Two", "change": "Trust newly scored 4", '
            '"implication": "They now publish licensing.", '
            '"recommendation": "Publish ours before Q4.", '
            '"classification": "transition_requirement", '
            '"decision_required": "Approve licensing page for Q4"}]}'
        )


@pytest.fixture
async def wm(tmp_path):
    from core.database import Database

    db = Database(str(tmp_path / "deck.db"))
    await db.initialize()
    yield WatchManager(db)
    await db.close()


async def _market(wm: WatchManager, c: str = "c1") -> dict[str, Any]:
    """Five brands: us, two rivals, one provisional, one never observed;
    a baseline snapshot, then three material moves."""
    dims = []
    for name, w in [
        ("Promos", 25),
        ("Payments", 25),
        ("Game library", 20),
        ("Support", 15),
        ("Trust", 15),
    ]:
        dims.append(
            await wm.upsert_dimension(
                name=name,
                company_id=c,
                weight_pct=w,
                subcriteria=[{"name": "a", "weight_pct": 50}, {"name": "b", "weight_pct": 50}],
            )
        )
    us = await wm.add_subject(name="OurBrand", company_id=c, url="https://u.example", is_self=True)
    r1 = await wm.add_subject(name="Rival One", company_id=c, url="https://r1.example")
    r2 = await wm.add_subject(name="Rival Two", company_id=c, url="https://r2.example")
    thin = await wm.add_subject(name="Thin Brand", company_id=c, url="https://t.example")
    await wm.add_subject(name="Ghost Brand", company_id=c, url="https://g.example")

    async def sc(subj, dim, score):
        await wm.add_evidence(
            company_id=c,
            subject_id=subj.subject_id,
            dimension_id=dim.dimension_id,
            claim=f"{subj.name} {dim.name}",
            source_url="https://x.example",
        )
        await wm.set_score(
            company_id=c, subject_id=subj.subject_id, dimension_id=dim.dimension_id, score=score
        )

    for d, v in zip(dims, [3, 4, 2, 5, 3], strict=True):
        await sc(us, d, v)
    for d, v in zip(dims, [5, 3, 4, 3, 4], strict=True):
        await sc(r1, d, v)
    for d, v in zip(dims[:4], [2, 2, 3, 2], strict=True):
        await sc(r2, d, v)
    await sc(thin, dims[3], 5)  # 15% of the model → provisional at 100.0
    await wm.take_snapshot(c, label="base")
    # Five material items: two score moves, one newly-scored dimension, and
    # the two rank changes those cause (us 2→1, Rival One 1→2).
    await sc(r1, dims[0], 2)  # material drop
    await sc(us, dims[2], 4)  # we improve
    await sc(r2, dims[4], 4)  # newly scored — always material
    return {"dims": dims, "us": us, "r1": r1, "r2": r2, "thin": thin}


def _slides(path: str) -> list[dict[str, Any]]:
    from pptx import Presentation

    out = []
    for sl in Presentation(path).slides:
        text = "\n".join(sh.text_frame.text for sh in sl.shapes if sh.has_text_frame)
        chart = next((sh.chart for sh in sl.shapes if sh.has_chart), None)
        table = next((sh.table for sh in sl.shapes if sh.has_table), None)
        pics = sum(1 for sh in sl.shapes if sh.shape_type == 13)  # PICTURE
        out.append(
            {
                "text": text,
                "chart": chart,
                "table": table,
                "pics": pics,
                "notes": sl.notes_slide.notes_text_frame.text,
            }
        )
    return out


def _slide(slides: list[dict[str, Any]], title_fragment: str) -> dict[str, Any]:
    for s in slides:
        if title_fragment.casefold() in s["text"].casefold():
            return s
    raise AssertionError(f"no slide contains {title_fragment!r}")


async def _deck(
    wm: WatchManager, tmp_path, router: Any = None, **params: Any
) -> tuple[Any, list[dict[str, Any]]]:
    from tools.watch.tools import WatchExecutiveDeckTool

    t = WatchExecutiveDeckTool()
    t._watch_manager = wm
    t._router = router
    res = await t.execute({"company_id": "c1", "path": str(tmp_path / "deck.pptx"), **params})
    assert res.success, res.error
    return res, _slides(res.data["path"])


class TestTheThreeRules:
    @pytest.mark.asyncio
    async def test_provisional_brand_is_listed_not_ranked(self, wm, tmp_path) -> None:
        """Thin Brand scores 100.0 on one 15%-weight dimension. That figure
        would lead the chart and read as 'market leader'. It must not."""
        await _market(wm)
        _res, slides = await _deck(wm, tmp_path)
        standings = _slide(slides, "Where the market stands")
        cats = list(standings["chart"].plots[0].categories)
        assert not any("Thin Brand" in c for c in cats)
        assert "Thin Brand" in standings["text"]  # named below the chart…
        assert "†" in standings["text"]  # …with the dagger and its reason
        assert "scored on 15% of total dimension weight" in standings["text"]

    @pytest.mark.asyncio
    async def test_unscored_is_blank_never_zero(self, wm, tmp_path) -> None:
        await _market(wm)
        _res, slides = await _deck(wm, tmp_path)
        # Chart: Ghost Brand has no scores — it gets no bar, not a 0 bar.
        standings = _slide(slides, "Where the market stands")
        vals = list(standings["chart"].plots[0].series[0].values)
        assert 0.0 not in vals
        assert not any("Ghost" in c for c in standings["chart"].plots[0].categories)
        assert "not yet observed" in standings["text"].casefold()
        assert "Ghost Brand" in standings["text"]
        # Heatmap: Rival Two's unscored cell is empty text, not "0".
        heat = _slide(slides, "Scores by dimension")
        tbl = heat["table"]
        header = [tbl.cell(0, c).text for c in range(len(tbl.columns))]
        rows = {
            tbl.cell(r, 0).text: [tbl.cell(r, c).text for c in range(len(tbl.columns))]
            for r in range(1, len(tbl.rows))
        }
        ghost = rows["Ghost Brand"]
        assert ghost[1] == "—" and all(v == "" for v in ghost[2:])
        thin = rows["Thin Brand"]
        assert thin[1] == "100.0†"
        assert thin[header.index("Support")] == "5"
        assert thin[header.index("Promos")] == ""

    @pytest.mark.asyncio
    async def test_our_bar_is_in_the_chart_and_ranked_order_holds(self, wm, tmp_path) -> None:
        await _market(wm)
        _res, slides = await _deck(wm, tmp_path)
        chart = _slide(slides, "Where the market stands")["chart"]
        cats = list(chart.plots[0].categories)
        vals = list(chart.plots[0].series[0].values)
        # python-pptx draws bar categories bottom-up, so #1 is last.
        assert cats[-1].startswith("OurBrand") and "(us)" in cats[-1]
        assert vals == sorted(vals)


class TestFactAndJudgementAreLabelled:
    @pytest.mark.asyncio
    async def test_without_a_model_the_summary_says_facts_only(self, wm, tmp_path) -> None:
        await _market(wm)
        res, slides = await _deck(wm, tmp_path)
        assert res.data["summary_source"] == "facts"
        summary = _slide(slides, "Executive summary")
        assert "Facts only" in summary["text"]
        # Reference-style summary: dimension columns + a "where we stand" strip.
        assert "OurBrand – #1" in summary["text"]
        assert "WHERE WE STAND" in summary["text"]
        # Baseline / change count is stated on the market-moves slide.
        assert "5 material" in _slide(slides, "Market moves")["text"] or any(
            "5 material" in s["text"] for s in slides
        )

    @pytest.mark.asyncio
    async def test_without_a_model_decisions_are_not_evaluated_not_none(self, wm, tmp_path) -> None:
        """'No decision required' and 'nobody judged it' must not look alike."""
        await _market(wm)
        _res, slides = await _deck(wm, tmp_path)
        decisions = _slide(slides, "Decisions required")
        assert "Not yet evaluated" in decisions["text"]
        assert "No board decision is required" not in decisions["text"]
        implications = _slide(slides, "Implications and recommendations")
        assert "no model was available" in implications["text"]

    @pytest.mark.asyncio
    async def test_with_a_model_summary_and_implications_are_labelled_as_judgement(
        self, wm, tmp_path
    ) -> None:
        await _market(wm)
        router = _DeckRouter()
        res, slides = await _deck(wm, tmp_path, router=router)
        assert res.data["summary_source"] == "model"
        assert res.data["judged_items"] == 2
        assert router.calls.count("judge") == 1 and router.calls.count("summary") == 1

        summary = _slide(slides, "Executive summary")
        assert "We lead, but Rival One is closing on promos." in summary["text"]
        assert "written by the model" in summary["text"]
        standings = _slide(slides, "OurBrand leads a thin field")
        assert "Amber is us – the lead is real but narrow." in standings["text"]

        implications = _slide(slides, "Implications and recommendations")
        tbl = implications["table"]
        classes = [tbl.cell(r, 4).text for r in range(1, len(tbl.rows))]
        assert classes == ["No-regret", "Transition requirement"]  # ordered by class
        assert "written by the model" in implications["text"]

        decisions = _slide(slides, "Decisions required")
        assert "Approve licensing page for Q4" in decisions["text"]
        assert "Rival One" not in decisions["text"]  # its decision was "none"

    @pytest.mark.asyncio
    async def test_a_broken_model_degrades_to_facts_and_says_so(self, wm, tmp_path) -> None:
        class Broken:
            async def complete(self, **_kw: Any) -> _Resp:
                return _Resp("not json at all")

        await _market(wm)
        res, slides = await _deck(wm, tmp_path, router=Broken())
        assert res.data["summary_source"] == "facts"
        assert "Facts only" in _slide(slides, "Executive summary")["text"]


class TestVersusTheLeader:
    @pytest.mark.asyncio
    async def test_compares_us_to_the_top_ranked_peer_dimension_by_dimension(
        self, wm, tmp_path
    ) -> None:
        await _market(wm)
        _res, slides = await _deck(wm, tmp_path)
        vs = _slide(slides, "Where we win, where we lose")
        # After the moves: us Promos 3 vs Rival One 2 → ahead; Trust 3 vs 4 → behind.
        assert "Promos (3 vs 2" in vs["text"]
        assert "Trust (3 vs 4" in vs["text"]
        assert "Rival One (#2)" in vs["text"]

    @pytest.mark.asyncio
    async def test_no_self_brand_says_so_instead_of_guessing(self, wm, tmp_path) -> None:
        d = await wm.upsert_dimension(
            name="Promos",
            company_id="c1",
            weight_pct=100,
            subcriteria=[{"name": "a", "weight_pct": 100}],
        )
        s = await wm.add_subject(name="Rival", company_id="c1", url="https://r.example")
        await wm.add_evidence(
            company_id="c1", subject_id=s.subject_id, dimension_id=d.dimension_id, claim="x"
        )
        await wm.set_score(
            company_id="c1", subject_id=s.subject_id, dimension_id=d.dimension_id, score=3
        )
        _res, slides = await _deck(wm, tmp_path)
        assert "No brand is marked as ours" in _slide(slides, "Where we win")["text"]


class TestFirstCycle:
    @pytest.mark.asyncio
    async def test_no_snapshot_means_baseline_not_zero_changes(self, wm, tmp_path) -> None:
        d = await wm.upsert_dimension(
            name="Promos",
            company_id="c1",
            weight_pct=100,
            subcriteria=[{"name": "a", "weight_pct": 100}],
        )
        s = await wm.add_subject(name="Rival", company_id="c1", url="https://r.example")
        await wm.add_evidence(
            company_id="c1", subject_id=s.subject_id, dimension_id=d.dimension_id, claim="x"
        )
        await wm.set_score(
            company_id="c1", subject_id=s.subject_id, dimension_id=d.dimension_id, score=3
        )
        res, slides = await _deck(wm, tmp_path)
        assert res.data["material_count"] == 0
        assert "First reporting cycle" in _slide(slides, "Baseline established")["text"]
        assert any("sets the baseline" in s["text"] for s in slides)


class TestEvidenceSlide:
    @pytest.mark.asyncio
    async def test_reports_coverage_and_the_no_penalty_rule(self, wm, tmp_path) -> None:
        await _market(wm)
        _res, slides = await _deck(wm, tmp_path)
        ev = _slide(slides, "How much of this is measured")
        # 5 brands × 5 dims = 25 pairs; scored: us 5 + r1 5 + r2 5 + thin 1 = 16.
        assert "16 of 25 brand × dimension pairs" in ev["text"]
        assert "64%" in ev["text"]
        assert "Gaps are absences of evidence, not weaknesses" in ev["text"]


class TestDeckShape:
    @pytest.mark.asyncio
    async def test_all_slides_within_bounds(self, wm, tmp_path) -> None:
        """12 core slides + one deep dive per real profiled competitor. The
        hallucinated brand in the router's profiles must not get a slide."""
        from pptx import Presentation

        await _market(wm)
        res, _ = await _deck(wm, tmp_path, router=_DeckRouter())
        prs = Presentation(res.data["path"])
        assert len(prs.slides) == 16
        W, H = prs.slide_width, prs.slide_height
        for sl in prs.slides:
            for sh in sl.shapes:
                assert sh.left >= 0 and sh.top >= 0
                assert sh.left + sh.width <= W and sh.top + sh.height <= H

    @pytest.mark.asyncio
    async def test_many_implications_paginate(self, wm, tmp_path) -> None:
        class Many:
            async def complete(self, *, messages, **_kw):
                if "executive-summary slide" in messages[0]["content"]:
                    return _Resp('{"headline": "h", "bullets": ["b"]}')
                items = ", ".join(
                    f'{{"subject": "S{i}", "change": "c", "implication": "i", '
                    f'"recommendation": "r", "classification": "monitor", '
                    f'"decision_required": "none"}}'
                    for i in range(12)
                )
                return _Resp(f'{{"items": [{items}]}}')

        from pptx import Presentation

        await _market(wm)
        res, slides = await _deck(wm, tmp_path, router=Many())
        titles = [s["text"] for s in slides if "Implications and recommendations" in s["text"]]
        assert len(titles) == 3  # 12 items, 5 per slide
        assert "(1/3)" in titles[0] and "(3/3)" in titles[2]
        # 12 core + dimension-leaders + 2 extra judgement pages, no profiles
        # (this router returns no profile section).
        assert len(Presentation(res.data["path"]).slides) == 16

    @pytest.mark.asyncio
    async def test_extension_is_forced_to_pptx(self, wm, tmp_path) -> None:
        await _market(wm)
        res, _ = await _deck(wm, tmp_path, path=str(tmp_path / "deck.pdf"))
        assert res.data["path"].endswith(".pptx")

    @pytest.mark.asyncio
    async def test_it_does_not_take_a_snapshot(self, wm, tmp_path) -> None:
        await _market(wm)
        before = len(await wm.list_snapshots("c1"))
        await _deck(wm, tmp_path)
        assert len(await wm.list_snapshots("c1")) == before


class TestBoardReportCanCarryTheDeck:
    @pytest.mark.asyncio
    async def test_deck_path_writes_both_from_one_period(self, wm, tmp_path) -> None:
        """Both artefacts are cut from the same diff, before the snapshot."""
        from tools.watch.tools import WatchBoardReportTool

        await _market(wm)
        t = WatchBoardReportTool()
        t._watch_manager = wm
        t._router = _DeckRouter()
        res = await t.execute(
            {
                "company_id": "c1",
                "path": str(tmp_path / "report.md"),
                "deck_path": str(tmp_path / "deck.pptx"),
                "take_snapshot": True,
            }
        )
        assert res.success
        assert res.data["material_count"] == 5
        assert res.data["deck_path"] == str(tmp_path / "deck.pptx")
        assert "deck_error" not in res.data
        slides = _slides(res.data["deck_path"])
        assert "5 material moves" in _slide(slides, "material moves this period")["text"]
        # …and the snapshot came after, so a follow-up diff is empty.
        assert (await wm.diff_since_snapshot("c1"))["material_count"] == 0

    @pytest.mark.asyncio
    async def test_report_to_disk_brings_the_deck_by_default(self, wm, tmp_path) -> None:
        """The deck is part of the pack, not an option to remember."""
        from tools.watch.tools import WatchBoardReportTool

        await _market(wm)
        t = WatchBoardReportTool()
        t._watch_manager = wm
        res = await t.execute(
            {"company_id": "c1", "path": str(tmp_path / "board-march.md"),
             "take_snapshot": False}
        )
        assert res.success
        assert res.data["deck_path"] == str(tmp_path / "board-march.pptx")
        assert (tmp_path / "board-march.pptx").exists()
        assert (tmp_path / "board-march.md").exists()

    @pytest.mark.asyncio
    async def test_deck_false_skips_it(self, wm, tmp_path) -> None:
        from tools.watch.tools import WatchBoardReportTool

        await _market(wm)
        t = WatchBoardReportTool()
        t._watch_manager = wm
        res = await t.execute(
            {"company_id": "c1", "path": str(tmp_path / "r.md"),
             "deck": False, "take_snapshot": False}
        )
        assert res.success and "deck_path" not in res.data
        assert not (tmp_path / "r.pptx").exists()

    @pytest.mark.asyncio
    async def test_inline_report_writes_nothing_so_no_deck(self, wm, tmp_path) -> None:
        """No path means the report is returned, not saved — same for the deck."""
        from tools.watch.tools import WatchBoardReportTool

        await _market(wm)
        t = WatchBoardReportTool()
        t._watch_manager = wm
        res = await t.execute({"company_id": "c1", "take_snapshot": False})
        assert res.success and "deck_path" not in res.data
        assert not list(tmp_path.glob("*.pptx"))


class TestRegistration:
    def test_tool_is_created_and_in_the_watch_group(self) -> None:
        from tools.watch.tools import create_watch_tools

        names = {t.name: t for t in create_watch_tools()}
        assert "watch_executive_deck" in names
        assert names["watch_executive_deck"].group == "watch"

    def test_analyze_ships_the_deck_by_default(self) -> None:
        from tools.watch.tools import WatchAnalyzeTool

        props = WatchAnalyzeTool().input_schema["properties"]
        assert "deck" in props
        assert "Default true" in props["deck"]["description"]
        assert "executive deck" in WatchAnalyzeTool().description


def _fake_png(path, w: int = 64, h: int = 40, rgb: tuple = (200, 60, 60)) -> str:
    """A real (tiny) PNG so python-pptx can measure it. Named .jpg or .png —
    pptx sniffs content, not extension."""
    import struct
    import zlib
    from pathlib import Path

    def chunk(typ: bytes, data: bytes) -> bytes:
        return (
            struct.pack(">I", len(data)) + typ + data
            + struct.pack(">I", zlib.crc32(typ + data) & 0xFFFFFFFF)
        )

    raw = b"".join(b"\x00" + bytes(rgb) * w for _ in range(h))
    png = (
        b"\x89PNG\r\n\x1a\n"
        + chunk(b"IHDR", struct.pack(">IIBBBBB", w, h, 8, 2, 0, 0, 0))
        + chunk(b"IDAT", zlib.compress(raw))
        + chunk(b"IEND", b"")
    )
    Path(path).parent.mkdir(parents=True, exist_ok=True)
    Path(path).write_bytes(png)
    return str(path)


class TestMarketFacingDeck:
    """The deck talks about the market: exec zones, competitor deep dives,
    storefront exhibits — and skips, never fakes, what it does not have."""

    @pytest.mark.asyncio
    async def test_exec_summary_reads_findings_threats_watch(self, wm, tmp_path) -> None:
        await _market(wm)
        _res, slides = await _deck(wm, tmp_path, router=_DeckRouter())
        summary = _slide(slides, "Executive summary")
        # The reference-deck shape: dimension columns with numbered observations,
        # then Recommendation / Where we stand / Decisions. The model's
        # by_dimension entries land in the columns; its recommendation in the strip.
        for zone in ("RECOMMENDATION", "WHERE WE STAND", "DECISIONS / NEXT STEPS"):
            assert zone in summary["text"]
        assert "Rival One leads promos with a 200% welcome bundle" in summary["text"]
        assert "Match the welcome bundle" in summary["text"]
        # Findings/threats/watch remain the fallback shape when no columns are given.

    @pytest.mark.asyncio
    async def test_competitor_profiles_render_observations_and_implications(
        self, wm, tmp_path
    ) -> None:
        await _market(wm)
        _res, slides = await _deck(wm, tmp_path, router=_DeckRouter())
        p1 = _slide(slides, "Rival One – closing fast on promotions")
        assert "COMPETITOR DEEP DIVE" in p1["text"]
        assert "two-tier welcome offer" in p1["text"]
        assert "Implications for us".upper() in p1["text"].upper()
        assert "Our promo calendar needs a counter" in p1["text"]
        assert "written by the model" in p1["text"]
        # The hallucinated brand was filtered before it reached a slide.
        with pytest.raises(AssertionError):
            _slide(slides, "hallucinated")

    @pytest.mark.asyncio
    async def test_facts_only_deck_has_no_profiles(self, wm, tmp_path) -> None:
        await _market(wm)
        _res, slides = await _deck(wm, tmp_path)  # no router
        assert not any("COMPETITOR DEEP DIVE" in s["text"] for s in slides)

    def test_exhibits_render_and_are_skipped_when_absent(self, tmp_path) -> None:
        from core.watch_deck import render_executive_deck

        card = {
            "generated_at": "2026-08-15T20:00:00+00:00",
            "weight_total_pct": 100.0,
            "dimensions": [{"name": "Promos", "weight_pct": 100.0}],
            "rows": [
                {
                    "name": "OurBrand", "is_self": True, "rank": 1,
                    "provisional": False,
                    "overall": {"normalized_pct": 80.0, "coverage_pct": 100.0},
                    "dimensions": {"Promos": {"score": 4}},
                },
                {
                    "name": "Rival One", "is_self": False, "rank": 2,
                    "provisional": False,
                    "overall": {"normalized_pct": 60.0, "coverage_pct": 100.0},
                    "dimensions": {"Promos": {"score": 3}},
                },
            ],
        }
        shots = {
            "OurBrand": [{
                "path": _fake_png(tmp_path / "us.jpg"),
                "url": "https://u.example", "observed_at": "2026-08-15T19:00:00",
            }],
            "Rival One": [{
                "path": _fake_png(tmp_path / "r1.jpg", rgb=(60, 60, 200)),
                "url": "https://r1.example", "observed_at": "2026-08-15T19:05:00",
            }],
        }
        with_shots = render_executive_deck(
            card, diff=None, judged=[], summary={"source": "facts"},
            gaps=[], evidence_count=2, screenshots=shots,
            path=tmp_path / "with.pptx",
        )
        slides = _slides(with_shots)
        ex = _slide(slides, "storefronts as a visitor sees them")
        assert ex["pics"] == 2
        assert "OurBrand" in ex["text"] and "(us)" in ex["text"]
        assert "captured 2026-08-15" in ex["text"]
        assert "state-verified network exit" in ex["text"]

        without = render_executive_deck(
            card, diff=None, judged=[], summary={"source": "facts"},
            gaps=[], evidence_count=2, path=tmp_path / "without.pptx",
        )
        assert not any(
            "storefronts as a visitor sees them" in s["text"].casefold()
            for s in _slides(without)
        )

    def test_profile_slide_carries_the_storefront_thumbnail(self, tmp_path) -> None:
        from core.watch_deck import render_executive_deck

        card = {
            "generated_at": "2026-08-15T20:00:00+00:00",
            "weight_total_pct": 100.0,
            "dimensions": [{"name": "Promos", "weight_pct": 100.0}],
            "rows": [{
                "name": "Rival One", "is_self": False, "rank": 1,
                "provisional": False,
                "overall": {"normalized_pct": 60.0, "coverage_pct": 100.0},
                "dimensions": {"Promos": {"score": 3}},
            }],
        }
        summary = {
            "source": "model",
            "profiles": [{
                "brand": "Rival One", "title": "the field's pace-setter",
                "observations": ["Daily bonus wheel on the homepage."],
                "implications": ["Sets the promo bar we are judged against."],
            }],
        }
        out = render_executive_deck(
            card, diff=None, judged=[], summary=summary, gaps=[],
            evidence_count=1,
            screenshots={"Rival One": [{
                "path": _fake_png(tmp_path / "r1.jpg"),
                "url": "https://r1.example", "observed_at": "2026-08-15",
            }]},
            path=tmp_path / "profile.pptx",
        )
        slides = _slides(out)
        p1 = _slide(slides, "the field's pace-setter")
        assert p1["pics"] == 1
        assert "Storefront capture" in p1["text"]

    def test_old_narrative_shape_still_renders(self, tmp_path) -> None:
        """A summary dict from before exec/profiles existed must render the
        bullets layout, not crash."""
        from core.watch_deck import render_executive_deck

        card = {
            "generated_at": "2026-08-15T20:00:00+00:00",
            "weight_total_pct": 100.0,
            "dimensions": [{"name": "Promos", "weight_pct": 100.0}],
            "rows": [],
        }
        out = render_executive_deck(
            card, diff=None, judged=[],
            summary={"headline": "Old shape", "bullets": ["Line one."],
                     "titles": {}, "commentary": {}, "next_steps": [],
                     "source": "facts"},
            gaps=[], evidence_count=0, path=tmp_path / "old.pptx",
        )
        slides = _slides(out)
        s = _slide(slides, "Old shape")
        assert "Line one." in s["text"]

    def test_eyebrows_cut_at_word_boundaries(self) -> None:
        """The cover once shipped reading '…FLORI'."""
        from core.watch_deck import _trim_words

        label = "Canonical 14-brand register · 12 weighted dimensions · Florida exits"
        cut = _trim_words(label, 60)
        assert len(cut) <= 60
        assert not cut.endswith("Flori")
        assert cut == cut.rstrip(" ·-–,;")
        assert _trim_words("short", 60) == "short"


class TestDimensionLeaders:
    """The per-dimension breakout: top observed score per battleground, our
    gap, and the honesty marks carried into the cells."""

    @pytest.mark.asyncio
    async def test_leaders_ties_daggers_and_gaps(self, wm, tmp_path) -> None:
        await _market(wm)
        _res, slides = await _deck(wm, tmp_path)
        sl = _slide(slides, "Who leads each dimension")
        tbl = sl["table"]
        rows = {
            tbl.cell(r, 0).text.split("  ·  ")[0]: [
                tbl.cell(r, c).text for c in range(len(tbl.columns))
            ]
            for r in range(1, len(tbl.rows))
        }
        # After the moves: us Promos 3, r1 dropped to 2, r2 has 2 → we lead.
        promos = rows["Promos"]
        assert "OurBrand (us)" in promos[1]
        assert promos[4] == "we lead"
        # Support: us 5, Thin Brand 5 (provisional overall) — a tie, and the
        # provisional brand carries its dagger into the leader cell.
        support = rows["Support"]
        assert "OurBrand (us)" in support[1]
        assert "Thin Brand †" in support[1]
        assert support[4] == "we co-lead"
        # Payments: us 4, r1 3, r2 2 → we lead outright at 4.
        payments = rows["Payments"]
        assert "OurBrand (us)" in payments[1]
        assert payments[2] == "4" and payments[3] == "4"
        # Trust: r1 4, r2 4, us 3 → behind by 1.
        trust = rows["Trust"]
        assert "OurBrand" not in trust[1]
        assert "behind" in trust[4] and "1" in trust[4]
        # Legend keeps the no-penalty rule on the slide.
        assert "never counted as a loss" in sl["text"]

    @pytest.mark.asyncio
    async def test_slide_sits_between_versus_and_deep_dives(
        self, wm, tmp_path
    ) -> None:
        await _market(wm)
        _res, slides = await _deck(wm, tmp_path, router=_DeckRouter())
        order = [
            i for i, s in enumerate(slides)
            if any(
                k in s["text"]
                for k in (
                    "VERSUS THE LEADER",
                    "Who leads each dimension",
                    "COMPETITOR DEEP DIVE",
                )
            )
        ]
        texts = [slides[i]["text"] for i in order]
        assert "VERSUS THE LEADER" in texts[0]
        assert "Who leads each dimension" in texts[1]
        assert "COMPETITOR DEEP DIVE" in texts[2]


class TestReferenceStyleDeck:
    """The customer's reference decks: every analytical slide carries a Key
    observations / Key implications panel; the offers on the table are a
    slide of their own; the summary is dimension columns + a decision strip."""

    @pytest.mark.asyncio
    async def test_reading_panels_and_offers_slide(self, wm, tmp_path) -> None:
        await _market(wm)
        _res, slides = await _deck(wm, tmp_path, router=_DeckRouter())
        standings = _slide(slides, "OurBrand leads a thin field")
        assert "Key observations" in standings["text"]
        assert "OurBrand leads; Rival One is one point back." in standings["text"]
        assert "Key implications" in standings["text"]
        offers = _slide(slides, "The offers on the table")
        # (table cells are not text frames; the panel and footnote are)
        assert "Rival One is the only brand stating a percentage welcome offer." in offers["text"]
        assert "Offer text as published" in offers["text"]
        # summary strip
        summary = _slide(slides, "Executive summary")
        assert "Approve a Q4 welcome bundle" in summary["text"]

    @pytest.mark.asyncio
    async def test_facts_only_deck_still_carries_computed_panels(self, wm, tmp_path) -> None:
        await _market(wm)
        _res, slides = await _deck(wm, tmp_path)
        standings = _slide(slides, "Where the market stands")
        assert "Key observations" in standings["text"]
        assert "leads the ranked field" in standings["text"]
